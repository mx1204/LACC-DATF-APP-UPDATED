import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import io
import re
import numpy as np
import traceback

# --- PAGE CONFIGURATION ---
st.set_page_config(page_title="Career Workshop Analytics Sandbox", layout="wide")

# --- CUSTOM CSS FOR "NOTEBOOK" FEEL ---
st.markdown("""
    <style>
    .block-container {padding-top: 2rem;}
    .stat-card {
        background-color: #f8f9fa;
        border-left: 0.3rem solid #4e73df;
        padding: 1rem;
        margin-bottom: 0.6rem;
        border-radius: 0.3rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .stat-title {
        font-size: clamp(0.85rem, 1.5vw, 1rem); 
        color: #555; 
        font-weight: bold;
    }
    .stat-value {
        font-size: clamp(1.2rem, 3vw, 1.8rem); 
        color: #2c3e50; 
        font-weight: bold;
    }
    h1, h2, h3 {color: #2c3e50;}
    .stTextArea textarea {font-family: 'Consolas', 'Courier New', monospace; background-color: #f4f4f4;}
    </style>
    """, unsafe_allow_html=True)

# --- HEADER ---
# Moved to main() for consistency across pages

# --- HELPER: SMART FILE LOADER ---
def load_data_smart(uploaded_file, file_type='attendance'):
    """
    Scans the first 10 rows to find the correct header row.
    For Excel files with multiple sheets, combines all sheets.
    Criteria:
    - Attendance: Look for 'Event Name' or 'Session Name'
    - Taxonomy: Look for 'Career Development Workshop Titles' or 'Event Name'
    """
    try:
        # Define keywords to search for in rows
        if file_type == 'attendance':
            keywords = ['event', 'session', 'title', 'date', 'time', 'status', 'simid', 'email', 'name', 'student', 'workshop', 'course', 'activity']
        else:
            keywords = ['career', 'workshop', 'event', 'category', 'title']
        
        # Handle Excel files with multiple sheets
        if uploaded_file.name.endswith('.xlsx') or uploaded_file.name.endswith('.xls'):
            try:
                uploaded_file.seek(0)
                excel_file = pd.ExcelFile(uploaded_file)
                all_dfs = []
                
                st.toast(f"📂 Processing {uploaded_file.name} ({len(excel_file.sheet_names)} sheets)...")
                
                for sheet_name in excel_file.sheet_names:
                    try:
                        # Read first few rows to find header
                        df_preview = pd.read_excel(excel_file, sheet_name=sheet_name, header=None, nrows=20)
                        
                        header_row_idx = 0
                        found_header = False
                        for idx, row in df_preview.iterrows():
                            row_str = row.astype(str).str.lower().str.strip().tolist()
                            if any(k in val for k in keywords for val in row_str):
                                header_row_idx = idx
                                found_header = True
                                break
                        
                        # Read full sheet
                        df_sheet = pd.read_excel(excel_file, sheet_name=sheet_name, header=header_row_idx)
                        
                        # Heuristic: Only add if it looks like a data table (has columns)
                        if len(df_sheet.columns) > 0:
                            all_dfs.append(df_sheet)
                            if found_header:
                                st.toast(f"✓ Loaded sheet '{sheet_name}' (Header: Row {header_row_idx+1})")
                            else:
                                st.toast(f"ℹ️ Loaded sheet '{sheet_name}' (Default Header)")
                        else:
                            st.toast(f"⚠️ Sheet '{sheet_name}' is empty.")
                            
                    except Exception as e:
                        st.error(f"Error reading sheet '{sheet_name}': {e}")

                if all_dfs:
                    df = pd.concat(all_dfs, ignore_index=True)
                    return df
                
                # Fallback: Hail Mary load
                st.warning(f"⚠️ Smart load failed for {uploaded_file.name}. Trying basic load...")
                uploaded_file.seek(0)
                return pd.read_excel(uploaded_file)
                
            except Exception as e:
                st.error(f"❌ Critical error loading Excel {uploaded_file.name}: {e}")
                return pd.DataFrame()
        
        # Handle CSV files
        else:
            # Try different encodings
            encodings = ['utf-8', 'ISO-8859-1', 'cp1252']
            for enc in encodings:
                try:
                    # Read first few rows without header to inspect
                    # Use sep=None to auto-detect delimiter (comma, semicolon, tab)
                    df_preview = pd.read_csv(uploaded_file, header=None, nrows=10, encoding=enc, sep=None, engine='python')
                    
                    header_row_idx = 0
                    found = False
                    
                    # Iterate through rows to find keywords
                    for idx, row in df_preview.iterrows():
                        row_str = row.astype(str).str.lower().str.strip().tolist()
                        if any(k in val for k in keywords for val in row_str):
                            header_row_idx = idx
                            found = True
                            break
                    
                    # Reset file pointer
                    uploaded_file.seek(0)
                    
                    # Read full file
                    df = pd.read_csv(uploaded_file, header=header_row_idx, encoding=enc, sep=None, engine='python')
                    
                    if not df.empty:
                        if found:
                            st.toast(f"✓ Loaded CSV '{uploaded_file.name}' (Header: Row {header_row_idx+1})")
                        else:
                            st.toast(f"ℹ️ Loaded CSV '{uploaded_file.name}' (Default Header)")
                        return df
                        
                except Exception as e:
                    continue # Try next encoding
            
            st.error(f"❌ Failed to load CSV {uploaded_file.name}. Checked encodings: {encodings}")
            return pd.DataFrame()
        
    except Exception as e:
        st.error(f"Error loading {uploaded_file.name}: {e}")
        return pd.DataFrame()

# --- STEP 2: BACKEND PROCESSING LOGIC ---

def run_data_ingestion():
    # --- STEP 1: DATA INGESTION ENGINE ---
    
    # Check if files are already loaded in session state to display them
    loaded_att = st.session_state.get('loaded_att_files', [])
    loaded_tax = st.session_state.get('loaded_tax_files', [])
    
    col_upload_1, col_upload_2 = st.columns(2)

    with col_upload_1:
        st.markdown("##### 1. Attendance Data")
        if loaded_att:
            st.success(f"✅ Loaded: {', '.join(loaded_att)}")
            with st.expander("Upload New / Replace"):
                attendance_files = st.file_uploader("Upload CSV/Excel", accept_multiple_files=True, key="att_files")
        else:
            attendance_files = st.file_uploader("Click to upload or drag and drop CSV/Excel files", accept_multiple_files=True, key="att_files")

    with col_upload_2:
        st.markdown("##### 2. Workshop Taxonomy")
        if loaded_tax:
             st.success(f"✅ Loaded: {', '.join(loaded_tax)}")
             with st.expander("Upload New / Replace"):
                taxonomy_files = st.file_uploader("Upload Excel/CSV", accept_multiple_files=True, key="tax_files")
        else:
            taxonomy_files = st.file_uploader("Click to upload or drag and drop Single Excel/CSV mapping file", accept_multiple_files=True, key="tax_files")

    # Determine if we should show the Process Button
    # Case A: New files uploaded -> Show Button
    # Case B: Already loaded -> Button not needed unless new files uploaded?
    # User said "maintain... unless user edit".
    # Logic: If new files are in uploader -> Process them.
    # If no new files but data exists -> Don't show process key, just show data.
    
    has_new_uploads = (attendance_files is not None and len(attendance_files) > 0) and (taxonomy_files is not None and len(taxonomy_files) > 0)
    

    # --- STEP 2: BACKEND PROCESSING LOGIC ---
    if 'data' not in st.session_state: st.session_state['data'] = None
    if 'review_mode' not in st.session_state: st.session_state['review_mode'] = False
    if 'taxonomy_categories' not in st.session_state: st.session_state['taxonomy_categories'] = []
    if 'debug_info' not in st.session_state: st.session_state['debug_info'] = {}

    # Show process button only when NEW files are uploaded
    if has_new_uploads:
        st.markdown("")
        process_button = st.button("Analyze Data", type="primary", use_container_width=True)
        st.markdown("")
    else:
        process_button = False

    # --- HOW IT WORKS SECTION ---
    # Show instructions when not in review mode (i.e., before data is processed or when returning to upload state)
    if not st.session_state.get('review_mode', False):
        st.markdown("### HOW IT WORKS")
        st.markdown("---")
        col1, col2, col3 = st.columns(3)
        with col1:
            st.info("**1. Upload Data**\n\nUpload your Attendance Logs (Box A) and Workshop Taxonomy (Box B) to begin.")
        with col2:
            st.info("**2. Review & Match**\n\nThe system auto-cleans data and matches events. You can manually resolve any uncategorized items.")
        with col3:
            st.info("**3. Analyze**\n\nRun specific analysis questions to generate interactive graphs, tables, and insights.")


    if has_new_uploads and process_button:
        # Save filenames to session state for persistence
        st.session_state['loaded_att_files'] = [f.name for f in attendance_files]
        st.session_state['loaded_tax_files'] = [f.name for f in taxonomy_files]
        
        with st.spinner('Processing and Merging Datasets...'):
            try:
                # 1. Load Attendance Data
                att_dfs = []
                for f in attendance_files:
                    df_temp = load_data_smart(f, file_type='attendance')
                    if not df_temp.empty:
                        att_dfs.append(df_temp)
                
                if att_dfs:
                    main_df = pd.concat(att_dfs, ignore_index=True)
                else:
                    st.error("❌ No valid attendance data loaded.")
                    st.stop()

                # 2. Load Taxonomy Data
                tax_dfs = []
                for f in taxonomy_files:
                    df_temp = load_data_smart(f, file_type='taxonomy')
                    if not df_temp.empty:
                        tax_dfs.append(df_temp)
                
                if tax_dfs:
                    tax_df = pd.concat(tax_dfs, ignore_index=True)
                else:
                    st.error("❌ No valid taxonomy data loaded.")
                    st.stop()

                # 3. Data Cleaning & Joining
                # --- ROBUST COLUMN HANDLING & MERGING ---
                # --- INDEPENDENT PROCESSING & HOLISTIC MATCHING ---
                
                # A. Standardize Columns (Strip whitespace only)
                main_df.columns = main_df.columns.str.strip()
                tax_df.columns = tax_df.columns.str.strip()

                # B. Identify Critical Columns
                # Attendance Data
                att_event_col = next((c for c in main_df.columns if c.lower() in ['event name', 'session name', 'title']), None)
                att_date_col = next((c for c in main_df.columns if c.lower() in ['date', 'attended date', 'attended_date', 'start date', 'event date']), None)
                
                if not att_event_col:
                    st.error(f"❌ Critical Error: Column 'Event Name' not found in Attendance Data. Available columns: {list(main_df.columns)}")
                    st.stop()

                # Taxonomy Data
                tax_title_col = next((c for c in tax_df.columns if c.lower() in ['career development workshop titles', 'workshop title', 'event name', 'title']), None)
                # Prioritize 'Sub-Category' over 'Category'
                tax_subcat_col = next((c for c in tax_df.columns if c.lower() in ['sub-category', 'sub category']), None)
                if not tax_subcat_col:
                    tax_subcat_col = next((c for c in tax_df.columns if c.lower() == 'category'), None)
                tax_date_col = next((c for c in tax_df.columns if c.lower() in ['date', 'event date']), None)
                tax_time_col = next((c for c in tax_df.columns if c.lower() in ['time', 'start time']), None)
                tax_day_col = next((c for c in tax_df.columns if c.lower() in ['day', 'day of week']), None)
                tax_trainer_col = next((c for c in tax_df.columns if c.lower() in ['trainer', 'instructor', 'facilitator', 'speaker']), None)
                
                if not tax_title_col:
                    tax_title_col = tax_df.columns[0]
                    st.warning(f"⚠️ Taxonomy Title column not found. Using '{tax_title_col}'.")

                # C. Robust Time Parsing Function
                def parse_time_range(time_str):
                    if not isinstance(time_str, str): return None
                    
                    # Standardize separators
                    t = time_str.lower().replace('.', ':').replace('to', '-').replace(' ', '')
                    
                    # Regex for formats like 3:30-6:30pm, 2-6pm, 09:00-17:00
                    import re
                    match = re.search(r'(\d{1,2}(?::\d{2})?)(?:[ap]m)?-(\d{1,2}(?::\d{2})?)([ap]m)?', t)
                    
                    if match:
                        start, end, end_period = match.groups()
                        
                        # Infer start period from end period if missing (e.g. 2-6pm -> 2pm start)
                        start_period = end_period if end_period else None
                        
                        # Helper to convert single time to 24h
                        def to_24h(val, period):
                            if ':' in val:
                                h, m = map(int, val.split(':'))
                            else:
                                h, m = int(val), 0
                            
                            if period == 'pm' and h != 12: h += 12
                            if period == 'am' and h == 12: h = 0
                            return f"{h:02d}:{m:02d}"

                        # If periods are ambiguous, standard business logic applies (9-5 is 9am-5pm)
                        # Simple heuristic: if end_period is pm, and start > end, start is am. 
                        # E.g. 10-2pm -> 10am-2pm. 2-4pm -> 2pm-4pm
                        
                        try:
                            # Determine meridiems
                            mp = 'pm' if 'pm' in end_period or 'pm' in t else ('am' if 'am' in t else None)
                            
                            # Default start period to same as end if not found? 
                            # Or if start > end (e.g. 10-2), start is am, end is pm.
                            
                            # Let's simplify: if period explicit in 3rd group, use it for end.
                            # For start, if not explicit, infer.
                            
                            # Just parse numbers first
                            def get_h(val):
                                if ':' in val: return int(val.split(':')[0])
                                return int(val)
                            
                            sh = get_h(start)
                            eh = get_h(end)
                            
                            sp = 'am'
                            ep = end_period.replace('am','').replace('pm','') if end_period else mp
                            
                            # Heuristic: 
                            # 9-5 -> 9am-5pm
                            # 2-4 -> 2pm-4pm (usually events are afternoon)
                            # 10-12 -> 10am-12pm
                            
                            if ep == 'pm':
                                if eh != 12: eh += 12
                                
                                # If start is larger than end (in 12h), it must be AM? No.
                                # If start < end (2-4), assume same period (pm)
                                if sh < 12 and sh < ( eh-12): # 2 < 4
                                    sh += 12
                                elif sh == 12:
                                    pass # 12pm
                                # else keep as am (9-5 -> 9 < 17, keep 9)
                            elif ep == 'am':
                               if eh == 12: eh = 0
                               if sh == 12: sh = 0
                            
                            # Reconstruct
                            def fmt(h, raw):
                                m = raw.split(':')[1] if ':' in raw else '00'
                                return f"{h:02d}:{m}"
                                
                            # This logic is getting complex. Let's trust the to_24h helper if we can set periods
                            # Let's use the provided examples as guide: "3.30-6.30PM, 2-6pm, and 9-5pm"
                            
                            # Case 1: 3.30-6.30PM -> Start PM, End PM.
                            # Case 2: 2-6pm -> Start PM, End PM (Attendance events unlikely 2am)
                            # Case 3: 9-5pm -> Start AM, End PM.
                            
                            # Revised Logic:
                            # If end contains PM:
                            #   If start hour >= 8 and start hour <= 11 -> AM
                            #   If start hour >= 1 and start hour <= 6 -> PM
                            #   If start hour == 12 -> PM
                            
                            is_pm = 'pm' in t
                            
                            final_start_h = sh
                            final_end_h = eh
                            
                            if is_pm:
                                if final_end_h < 12: final_end_h += 12 # 6pm -> 18
                                if final_end_h == 12: pass # 12pm -> 12
                                
                                # Deduce start
                                if 8 <= final_start_h <= 11: # 9-5pm -> 9am
                                    pass 
                                elif 1 <= final_start_h <= 6: # 2-6pm -> 2pm -> 14
                                    final_start_h += 12
                                elif final_start_h == 12: # 12-2pm -> 12pm
                                    pass
                                    
                            def fmt2(h, raw):
                                 m = raw.split(':')[1] if ':' in raw else '00'
                                 return f"{h:02d}:{m}"

                            return f"{fmt2(final_start_h, start)}-{fmt2(final_end_h, end)}"

                        except:
                            return time_str # Fallback
                    return time_str

                # Apply Time Parsing to Taxonomy
                if tax_time_col:
                    tax_df['Formatted_Time'] = tax_df[tax_time_col].apply(parse_time_range)
                else:
                    tax_df['Formatted_Time'] = None

                # D. Pre-process Dates for Strict Matching
                # We strictly need valid dates to match.
                main_df['match_date'] = pd.to_datetime(main_df[att_date_col], dayfirst=True, errors='coerce').dt.date if att_date_col else None
                tax_df['match_date'] = pd.to_datetime(tax_df[tax_date_col], dayfirst=True, errors='coerce').dt.date if tax_date_col else None

                # E. Token Overlap Logic Helper
                def get_tokens(text):
                    if not isinstance(text, str): return set()
                    # Remove punctuation, lower case
                    clean = re.sub(r'[^\w\s]', '', text.lower())
                    return set(clean.split())

                # Prepare Taxonomy Lookup for efficiency
                # Structure: { match_date: [ {tokens: set, subcat: str, time: str, title: str} ] }
                tax_lookup = {}
                for idx, row in tax_df.iterrows():
                    d = row['match_date']
                    if pd.isna(d): continue
                    
                    if d not in tax_lookup: tax_lookup[d] = []
                    
                    tax_lookup[d].append({
                        'tokens': get_tokens(row[tax_title_col]),
                        'subcat': row[tax_subcat_col] if tax_subcat_col else 'Uncategorized',
                        'time': row['Formatted_Time'],
                        'day': row[tax_day_col] if tax_day_col else None, # Store Day
                        'trainer': row[tax_trainer_col] if tax_trainer_col else None, # Store Trainer
                        'original_title': row[tax_title_col]
                    })

                # F. Execution of Matching
                st.info("🔍 Matching Attendance with Taxonomy (Strict Date + Title Token Overlap)...")

                # REFINED LOGIC: Create Title -> Trainer Map for Exact Matching (Independent of Date)
                # This ensures that if we know the Trainer for a specific Workshop Title, we use it 
                # even if the dates don't perfectly align or are missing.
                title_trainer_map = {}
                if tax_trainer_col and tax_title_col:
                    # Drop rows where Title or Trainer is missing
                    temp_map = tax_df.dropna(subset=[tax_title_col, tax_trainer_col])
                    # Create dictionary (Title -> Trainer)
                    # Use last occurrence or mode? Standard dict(zip) effectively uses last content found.
                    title_trainer_map = dict(zip(temp_map[tax_title_col], temp_map[tax_trainer_col]))

                def match_record(row):
                    # Default
                    res = {'Sub-Category': 'Uncategorized', 'Matched_Time': None, 'Matched_Day': None, 'Matched_Trainer': None}
                    
                    # 0. Refined Logic: Exact Title Lookup just for Trainer
                    # If we recognize the title, pre-fill the Trainer
                    if row[att_event_col] in title_trainer_map:
                        res['Matched_Trainer'] = title_trainer_map[row[att_event_col]]

                    # Get record details
                    r_tokens = get_tokens(row[att_event_col])
                    
                    best_match = None
                    max_overlap = 0
                    
                    # Title-Only Match across ALL taxonomy entries
                    # Collect all candidates from all dates
                    all_candidates = []
                    for date_candidates in tax_lookup.values():
                        all_candidates.extend(date_candidates)
                    
                    # Token-based matching
                    for candidate in all_candidates:
                        c_tokens = candidate['tokens']
                        if not c_tokens: continue
                        
                        intersection = r_tokens.intersection(c_tokens)
                        overlap_score = len(intersection) / len(c_tokens) if len(c_tokens) > 0 else 0
                        
                        # Use 70% threshold for quality matching
                        if overlap_score >= 0.7 and overlap_score > max_overlap:
                            max_overlap = overlap_score
                            best_match = candidate
                    
                    # Fuzzy fallback if token matching didn't work
                    if not best_match and all_candidates:
                        import difflib
                        all_titles = [c['original_title'] for c in all_candidates]
                        matches = difflib.get_close_matches(row[att_event_col], all_titles, n=1, cutoff=0.7)
                        
                        if matches:
                            match_title = matches[0]
                            for c in all_candidates:
                                if c['original_title'] == match_title:
                                    best_match = c
                                    break
                    
                    # Apply result
                    if best_match:
                        res['Sub-Category'] = best_match['subcat']
                        res['Matched_Time'] = best_match['time']
                        res['Matched_Day'] = best_match['day']
                        
                        # Only overwrite Trainer if the Date-Specific match has one
                        if best_match['trainer']:
                            res['Matched_Trainer'] = best_match['trainer']
                        
                    return pd.Series(res)

                # Apply matching
                matched_cols = main_df.apply(match_record, axis=1)
                # Ensure matched_cols has specific columns in order
                matched_cols = matched_cols[['Sub-Category', 'Matched_Time', 'Matched_Day', 'Matched_Trainer']]
                
                # Save available categories for Manual Resolution
                if tax_subcat_col:
                    # specific categories + [DELETE]
                    cats = sorted(tax_df[tax_subcat_col].dropna().astype(str).str.lower().unique().tolist())
                    # Add [DELETE] as the first option
                    cats.insert(0, '[DELETE]')
                    st.session_state['taxonomy_categories'] = cats
                else:
                    st.session_state['taxonomy_categories'] = ['[DELETE]']
                
                merged_df = pd.concat([main_df, matched_cols], axis=1)
                
                # Rename Matched_Day to Workshop Timing_Day
                if 'Matched_Day' in merged_df.columns:
                     merged_df.rename(columns={'Matched_Day': 'Workshop Timing_Day'}, inplace=True)
                
                
                # FALLBACK: Derive Day from Attended Date if missing
                if 'Workshop Timing_Day' not in merged_df.columns and 'Attended Date' in merged_df.columns:
                     merged_df['Workshop Timing_Day'] = pd.to_datetime(merged_df['Attended Date'], errors='coerce').dt.day_name()
                # Also fill NaNs in existing column
                elif 'Workshop Timing_Day' in merged_df.columns and 'Attended Date' in merged_df.columns:
                     merged_df['Workshop Timing_Day'] = merged_df['Workshop Timing_Day'].fillna(
                         pd.to_datetime(merged_df['Attended Date'], errors='coerce').dt.day_name()
                     )

                # Cleanup
                if 'match_date' in merged_df.columns:
                    merged_df.drop(columns=['match_date'], inplace=True)
                if 'temp_match_date' in merged_df.columns:
                     merged_df.drop(columns=['temp_match_date'], inplace=True)
                     
                # Save Categories for Manual Resolution
                if tax_subcat_col:
                    st.session_state['taxonomy_categories'] = ['[DELETE]'] + sorted(tax_df[tax_subcat_col].dropna().astype(str).str.lower().unique().tolist())
                else:
                    st.session_state['taxonomy_categories'] = []
                
                # Debug Stats
                st.session_state['debug_info'] = {
                    'match_attempted': len(merged_df),
                    'matched': (merged_df['Sub-Category'] != 'Uncategorized').sum()
                }

                # 4. Derive Attributes for Analysis
                # Map standard analysis columns but do NOT rename original columns to lose data.
                # Instead, we COPY relevant data to standardized columns if missing.
                
                # Map of Standard Name -> Potential Original Columns
                col_map = {
                    'Event Name': [att_event_col],
                    'Attended Date': [att_date_col],
                    'Attendance Status': ['attendance status', 'status', 'attendance', 'registration status', 'participant status', 'rsvp status'],
                    'Registered Date': ['registered date', 'registration date', 'reg date', 'created at', 'timestamp'],
                    'Expected Grad Term': ['expected grad term', 'grad term', 'graduation date', 'expected graduation', 'grad date'],
                    'Citizenship': ['citizenship', 'citizen'],
                    'Nationality': ['nationality', 'nation'],
                    'SIMID': ['simid', 'student id', 'id', 'sim id', 'admin number'],
                    'University Program': ['university program', 'university', 'institution', 'partner university', 'school'],
                    'Program': ['program', 'major', 'course', 'degree', 'study'],
                    'Original_Time': ['time', 'attended time', 'attended_time', 'start time', 'check in time'], # Keep original time separate
                    'Original_Trainer': ['trainer', 'instructor', 'facilitator', 'speaker', 'presenter'],
                    'Student Name': ['student name', 'name', 'full name', 'participant name', 'attendee name']
                }
                
                # Ensure standard columns exist for analysis code
                for std_col, candidates in col_map.items():
                    if std_col not in merged_df.columns:
                        # Find first match
                        found = next((c for c in merged_df.columns if c.lower() in [x.lower() for x in candidates]), None)
                        if found:
                            merged_df[std_col] = merged_df[found] # Create copy for analysis
                        else:
                            merged_df[std_col] = pd.NA
                
                # Consolidate Time
                # Logic: If we found a Matched_Time (from Taxonomy), use it as the analysis 'Time'.
                # Else, fall back to Original_Time.
                # Add new standardized 'Time' column for analysis.
                merged_df['Time'] = merged_df['Matched_Time'].combine_first(merged_df['Original_Time'])
                
                # Consolidate Trainer
                # Prioritize Taxonomy
                merged_df['Trainer'] = merged_df['Matched_Trainer'].combine_first(merged_df['Original_Trainer'])
                
                if 'Matched_Time' in merged_df.columns:
                     st.toast("✓ Synced 'Time' with Taxonomy schedule where dates matched.")
                
                # --- NORMALIZE ATTENDANCE STATUS ---
                # Ensure 'Attendance Status' has standard values for analysis
                if 'Attendance Status' in merged_df.columns:
                    def normalize_status(val):
                        s = str(val).lower().strip()
                        if s in ['attended', 'present', 'checked in', 'yes', 'completed', 'participated', 'show']:
                            return 'Attended'
                        elif s in ['absent', 'no show', 'cancelled', 'no', 'registered']:
                            return 'Absent'
                        return val # Keep original if unknown
                    
                    merged_df['Attendance Status'] = merged_df['Attendance Status'].apply(normalize_status)
                else:
                    # If column missing, assume everyone attended (since it's an attendance log)
                    merged_df['Attendance Status'] = 'Attended'
                    st.warning("⚠️ 'Attendance Status' column not found. Assuming all records are 'Attended'.")

                # Ensure required columns exist
                for col in ['Attended Date', 'Registered Date', 'Expected Grad Term', 'Student Name']:
                    if col not in merged_df.columns:
                        merged_df[col] = pd.NA

                if 'Time' not in merged_df.columns:
                    merged_df['Time'] = None
                    
                # Derive 'Workshop Timing_Hours'
                # Format: '15:30-18:30' -> 15 (integer)
                # Also handle raw strings like '7-9pm' -> 19 using AM/PM context
                def get_start_hour(time_val):
                    if not isinstance(time_val, str): return pd.NA
                    s_val = time_val.lower().strip()
                    try:
                        # Get start part
                        start_part = s_val.split('-')[0].strip() # '15:30' or '7' or '7pm'
                        
                        # Extract hour number
                        if ':' in start_part:
                            h = int(start_part.split(':')[0])
                        else:
                            # Extract digits only from start part
                            digits = "".join(filter(str.isdigit, start_part))
                            if not digits: return pd.NA
                            h = int(digits)
                        
                        # Heuristic for AM/PM if not in 24h format (i.e., h <= 12)
                        # If 'pm' is in the full string, we might need to adjust
                        if 'pm' in s_val and h <= 12:
                            # Cases:
                            # 1. '2pm-...' -> Explicit PM in start matching
                            if 'pm' in start_part:
                                if h < 12: h += 12
                            # 2. '9am-5pm' -> Explicit AM in start matching
                            elif 'am' in start_part:
                                if h == 12: h = 0
                            # 3. '2-6pm' -> Implicit PM (range end has PM)
                            else:
                                # Start < 12. 
                                # If start is 8, 9, 10, 11 -> Likely AM (9-5pm)
                                # If start is 1, 2, 3, 4, 5, 6, 7 -> Likely PM (2-6pm)
                                if 1 <= h <= 7:
                                    h += 12
                                # 12 is 12pm, keep it. 8-11 keep it (AM).
                                
                        elif 'am' in s_val and h == 12:
                             h = 0 # 12am -> 0
                             
                        return h
                    except:
                        return pd.NA
                        
                merged_df['Workshop Timing_Hours'] = merged_df['Time'].apply(get_start_hour)

                # --- ROBUST DATE PARSING ---
                def parse_dates(series):
                    # Try default first
                    d = pd.to_datetime(series, errors='coerce')
                    # If too many NaTs, try dayfirst
                    if d.isna().sum() > 0.5 * len(d):
                        d = pd.to_datetime(series, errors='coerce', dayfirst=True)
                    # If still too many NaTs, try mixed format
                    if d.isna().sum() > 0.5 * len(d):
                        d = pd.to_datetime(series, errors='coerce', format='mixed')
                    return d

                merged_df['Attended Date'] = parse_dates(merged_df['Attended Date'])
                merged_df['Registered Date'] = parse_dates(merged_df['Registered Date'])
                
                # Create a backup of the original string for the final export
                if 'Expected Grad Term' in merged_df.columns:
                    merged_df['Original_Grad_Term'] = merged_df['Expected Grad Term']
                
                # --- PARSE EXPECTED GRAD TERM INTO YYYY-MM FORMAT ---
                def parse_grad_term(grad_term_str):
                    """
                    Parse Expected Grad Term text into year and month.
                    Examples:
                    - "2025-26 (Aug-Jul)" → Year: 2026, Month: Jul (07)
                    - "2025 Semester (Jan-Jun)" → Year: 2025, Month: Jun (06)
                    - "2024-25 (Sep-Aug)" → Year: 2025, Month: Aug (08)
                    """
                    if not isinstance(grad_term_str, str):
                        return pd.NA, pd.NA, pd.NA
                    
                    # Month mapping
                    month_map = {
                        'jan': ('01', 'January'), 'feb': ('02', 'February'), 'mar': ('03', 'March'),
                        'apr': ('04', 'April'), 'may': ('05', 'May'), 'jun': ('06', 'June'),
                        'jul': ('07', 'July'), 'aug': ('08', 'August'), 'sep': ('09', 'September'),
                        'oct': ('10', 'October'), 'nov': ('11', 'November'), 'dec': ('12', 'December')
                    }
                    
                    grad_year = pd.NA
                    grad_month_num = pd.NA
                    grad_month_name = pd.NA
                    
                    try:
                        # Extract month from parentheses - take the END month (second month)
                        month_match = re.search(r'\(.*?-\s*([A-Za-z]+)\s*\)', grad_term_str)
                        if month_match:
                            end_month = month_match.group(1).lower()[:3]  # Get first 3 letters
                            if end_month in month_map:
                                grad_month_num, grad_month_name = month_map[end_month]
                        
                        # Extract year - look for year patterns
                        # Case A: Range format like "2025-26" → take second year
                        year_range_match = re.search(r'(\d{4})-(\d{2})', grad_term_str)
                        if year_range_match:
                            first_year = year_range_match.group(1)[:2]  # "20"
                            second_year = year_range_match.group(2)  # "26"
                            grad_year = int(first_year + second_year)  # "2026"
                        else:
                            # Case B: Single year like "2025"
                            year_match = re.search(r'(\d{4})', grad_term_str)
                            if year_match:
                                grad_year = int(year_match.group(1))
                        
                    except:
                        pass  # Return NAs on any error
                    
                    return grad_year, grad_month_name, grad_month_num
                
                # Apply parsing if Original_Grad_Term exists (NO 'Workshop Timing_' prefix)
                if 'Original_Grad_Term' in merged_df.columns:
                    parsed_grad = merged_df['Original_Grad_Term'].apply(parse_grad_term)
                    merged_df['Grad_Year'] = parsed_grad.apply(lambda x: x[0])
                    merged_df['Grad_Month'] = parsed_grad.apply(lambda x: x[1])
                    grad_month_num = parsed_grad.apply(lambda x: x[2])
                    
                    # Create YYYY-MM format
                    merged_df['Grad_YYYY-MM'] = merged_df.apply(
                        lambda row: f"{int(row['Grad_Year'])}-{grad_month_num[row.name]}" 
                        if pd.notna(row['Grad_Year']) and pd.notna(grad_month_num[row.name])
                        else pd.NA, axis=1
                    )
                else:
                    merged_df['Grad_Year'] = pd.NA
                    merged_df['Grad_Month'] = pd.NA
                    merged_df['Grad_YYYY-MM'] = pd.NA
                
                # DO NOT parse Expected Grad Term to datetime - keep original string values
                # merged_df['Expected Grad Term'] = parse_dates(merged_df['Expected Grad Term'])  # REMOVED
                
                # Year (Strict: Only from Attended Date)
                merged_df['Workshop Timing_Year'] = merged_df['Attended Date'].dt.year
                
                # Month (Month name, e.g., "January", "February")
                merged_df['Workshop Timing_Month'] = merged_df['Attended Date'].dt.month_name()
                
                # DayNumber (Day of month, 1-31)
                merged_df['Workshop Timing_DayNumber'] = merged_df['Attended Date'].dt.day

                # Lead Time (Days) - Reconstruct attended date from components and subtract from Registered Date
                # Create a proper datetime from Workshop Timing components
                merged_df['_temp_attended_date'] = pd.to_datetime(
                    merged_df[['Workshop Timing_Year', 'Workshop Timing_Month', 'Workshop Timing_DayNumber']].rename(
                        columns={'Workshop Timing_Year': 'year', 'Workshop Timing_Month': 'month', 'Workshop Timing_DayNumber': 'day'}
                    ),
                    errors='coerce'
                )
                merged_df['Lead_Days'] = (merged_df['_temp_attended_date'] - merged_df['Registered Date']).dt.days
                merged_df.drop(columns=['_temp_attended_date'], inplace=True)
                
                # Student Type (Local/Intl)
                # Logic: Use Citizenship column. If NaN or empty, treat as 'Local'
                if 'Citizenship' in merged_df.columns:
                    def get_student_type(c):
                        # If NaN or empty, treat as Local
                        if pd.isna(c) or str(c).strip() == '':
                            return 'Local'
                        c_str = str(c).lower()
                        if 'singapore' in c_str or 'pr' in c_str:
                            return 'Local'
                        return 'International'
                    merged_df['Student_Type'] = merged_df['Citizenship'].apply(get_student_type)
                else:
                    # Fallback: If Citizenship column doesn't exist, default to 'Local'
                    merged_df['Student_Type'] = 'Local'

                # --- PROGRAM CLASSIFICATION ---
                # Categorize academic majors into broad categories with priority-based matching
                def categorize_program(program):
                    """
                    Categorizes academic programs into broad categories.
                    Priority-based matching: Foundation/Grad → Tech → Social → Business → Other
                    """
                    if pd.isna(program):
                        return 'not specified'
                    
                    program_lower = str(program).lower()
                    
                    # Priority 1: Non-Degree/Advanced Programs (Graduate, Diploma, Foundation)
                    foundation_keywords = [
                        'graduate diploma', 'master of science', 'diploma', 'foundation programme',
                        'foundation program', 'pre-university', 'postgraduate', 'graduate certificate',
                        'advanced diploma', 'international foundation'
                    ]
                    if any(keyword in program_lower for keyword in foundation_keywords):
                        return 'graduate, diploma, & foundation'
                    
                    # Priority 2: Technical Fields (Computing, IT, Engineering)
                    tech_keywords = [
                        'data science', 'computer science', 'cyber security', 'machine learning',
                        'artificial intelligence', 'information technology', 'software engineering',
                        'computer engineering', 'computing', 'it', 'cybersecurity', 'ai', 'ml',
                        'data analytics', 'information systems', 'network', 'programming'
                    ]
                    if any(keyword in program_lower for keyword in tech_keywords):
                        return 'computing & it'
                    
                    # Priority 3: Social Sciences & Arts
                    social_keywords = [
                        'psychology', 'communication', 'sociology', 'arts', 'humanities',
                        'social science', 'political science', 'history', 'philosophy',
                        'english', 'literature', 'media', 'journalism', 'design'
                    ]
                    if any(keyword in program_lower for keyword in social_keywords):
                        return 'social sciences & arts'
                    
                    # Priority 4: Business & Finance
                    business_keywords = [
                        'business', 'management', 'accounting', 'finance', 'economics',
                        'commerce', 'marketing', 'entrepreneurship', 'business administration',
                        'mba', 'financial', 'banking', 'international business'
                    ]
                    if any(keyword in program_lower for keyword in business_keywords):
                        return 'business & finance'
                    
                    # Default: Other
                    return 'other'
                
                # Apply program classification
                if 'Program' in merged_df.columns:
                    merged_df['Program Classification'] = merged_df['Program'].apply(categorize_program)
                else:
                    merged_df['Program Classification'] = 'not specified'



                # --- FIX SUB-CATEGORY TYPO ---
                # Correct "per-onboarding" to "pre-onboarding"
                if 'Sub-Category' in merged_df.columns:
                    merged_df['Sub-Category'] = merged_df['Sub-Category'].astype(str).str.replace(
                        'per-onboarding', 'pre-onboarding', case=False, regex=False
                    )

                # --- NORMALIZE ALL TEXT COLUMNS TO LOWERCASE ---
                # This ensures case-insensitive operations throughout the app
                text_columns = merged_df.select_dtypes(include=['object']).columns
                for col in text_columns:
                    # Skip columns that should preserve case (like SIMID)
                    if col not in ['SIMID', 'match_key', 'Time']:
                        try:
                            merged_df[col] = merged_df[col].astype(str).str.lower()
                        except:
                            pass  # Skip if conversion fails


                # --- BACKFILL WORKSHOP TIMING ATTRIBUTES FOR ABSENT RECORDS ---
                # For records with missing Workshop Timing data (typically Absent status),
                # find the closest record with the same Event Name and copy timing attributes
                st.info("🔄 Backfilling Workshop Timing attributes for records with missing data...")
                
                # Identify Workshop Timing columns
                workshop_timing_cols = [col for col in merged_df.columns if col.startswith('Workshop Timing_')]
                
                if workshop_timing_cols:
                    # Find records with missing Workshop Timing data
                    missing_mask = merged_df[workshop_timing_cols].isna().any(axis=1)
                    records_with_missing = merged_df[missing_mask].index.tolist()
                    
                    backfilled_count = 0
                    
                    for idx in records_with_missing:
                        # Get the event name for this record
                        event_name = merged_df.loc[idx, 'Event Name']
                        
                        if pd.isna(event_name):
                            continue
                        
                        # Search for nearby records with same event name and complete timing data
                        # Search range: 10 rows above and below
                        search_range = 10
                        start_idx = max(0, idx - search_range)
                        end_idx = min(len(merged_df), idx + search_range + 1)
                        
                        # Get nearby records
                        nearby_records = merged_df.iloc[start_idx:end_idx]
                        
                        # Filter for same event name and complete timing data
                        same_event = nearby_records[nearby_records['Event Name'] == event_name]
                        complete_timing = same_event[~same_event[workshop_timing_cols].isna().any(axis=1)]
                        
                        if not complete_timing.empty:
                            # Find the closest record (by index distance)
                            # Convert to Series to use abs() method
                            distances = pd.Series(complete_timing.index - idx).abs()
                            closest_idx = complete_timing.index[distances.idxmin()]
                            
                            # Copy all Workshop Timing attributes from closest record
                            for col in workshop_timing_cols:
                                merged_df.loc[idx, col] = merged_df.loc[closest_idx, col]
                            
                            backfilled_count += 1
                    
                    if backfilled_count > 0:
                        st.success(f"✓ Backfilled Workshop Timing data for {backfilled_count} records")
                    else:
                        st.info("ℹ️ No records needed backfilling")

                # --- ADD UNIQUE ID COLUMN ---
                # Generate sequential ID starting from 1
                merged_df.insert(0, 'ID', range(1, len(merged_df) + 1))

                # --- CLEANUP INTERMEDIATE COLUMNS ---
                # Remove columns requested to be excluded from final dataset, if they are not strictly needed for global state
                # 'Original_Time' is now merged into 'Time', so it's redundant.
                # 'Lead_Days' is recalculated in Q10, so we can remove it from global state to keep it clean.
                # 'Session Name' and 'Date' are removed as requested by user.
                # 'Registered Date' is KEPT in the final dataset.
                # We keep 'Attended Date' components (Year, Month, DayNumber) instead of raw 'Date'.
                cols_to_drop = ['Original_Time', 'Lead_Days', 'Matched_Time', 'Session Name', 'Date']
                merged_df.drop(columns=[c for c in cols_to_drop if c in merged_df.columns], inplace=True)

                # Store in Session State & Enter Review Mode
                st.session_state['data'] = merged_df
                st.session_state['review_mode'] = True
                
                # --- AUTO NAVIGATE IF REQUESTED ---
                # But we can't auto-navigate here because we are in a button callback?
                # Actually, rerun will reload the app. If we set current_page to 2? NO.
                # The user said: "when click the last button at page 1, auto go to page 2"
                # The button is "Finish Data Prep & Go to Analysis".
                # That button is BELOW. This block is for "Analyze Data" (initial processing).
                
                st.rerun()

            except Exception as e:
                st.error(f"Error processing files: {e}")
                st.stop()

    # PHASE 2: REVIEW & MANUAL RESOLUTION
    if st.session_state['review_mode'] and st.session_state['data'] is not None:
        merged_df = st.session_state['data']
        
        st.info("📊 **Data Processing Complete!** Please review the matching statistics below.")
        
        # --- STATISTICS ---
        col_stat1, col_stat2, col_stat3 = st.columns(3)
        with col_stat1:
            st.metric("Total Events", len(merged_df))
        with col_stat2:
            matched_count = merged_df['Sub-Category'].ne('uncategorized').sum()
            match_pct = matched_count/len(merged_df)*100
            st.metric("Categorized Events", f"{matched_count} ({match_pct:.1f}%)")
        with col_stat3:
            st.write("") # Spacer
            # Only enable download if all items are resolved
            current_uncategorized = (merged_df['Sub-Category'] == 'uncategorized').sum()
            
            if current_uncategorized == 0:
                csv = merged_df.to_csv(index=False).encode('utf-8')
                st.download_button(
                    label="📥 Download Merged Data",
                    data=csv,
                    file_name="merged_workshop_data.csv",
                    mime="text/csv"
                )
            else:
                st.caption(f"Resolve {current_uncategorized} items to download.")
        
        # --- MANUAL RESOLUTION TOOL ---
        uncategorized_count = (merged_df['Sub-Category'] == 'uncategorized').sum()
        
        if uncategorized_count > 0:
            st.warning(f"⚠️ **Action Needed**: {uncategorized_count} records remain uncategorized.")
            
            with st.expander("🛠️ Manual Resolution Tool - Assign or Delete Uncategorized Events", expanded=True):
                st.markdown("**All Uncategorized Events:**")
                st.caption("For each unique event below, either assign a Sub-Category or mark for deletion.")
                
                # Get all unique uncategorized event names with their counts
                uncategorized_events = merged_df[merged_df['Sub-Category'] == 'uncategorized']['Event Name'].value_counts()
                
                # Get available sub-categories from taxonomy
                available_categories = st.session_state.get('taxonomy_categories', [])
                
                # Create a form for batch processing
                with st.form("manual_resolution_form"):
                    assignments = {}
                    
                    # Display each uncategorized event with a dropdown
                    for idx, (event_name, count) in enumerate(uncategorized_events.items()):
                        col1, col2 = st.columns([3, 1])
                        
                        with col1:
                            st.markdown(f"**{idx+1}. {event_name}**")
                            st.caption(f"Appears {count} time(s)")
                        
                        with col2:
                            selected_category = st.selectbox(
                                "Action",
                                options=available_categories,
                                index=0,  # Default to [DELETE]
                                key=f"cat_{idx}",
                                label_visibility="collapsed"
                            )
                            assignments[event_name] = selected_category
                        
                        if idx < len(uncategorized_events) - 1:
                            st.markdown("---")
                    
                    # Submit button
                    st.markdown("---")
                    submitted = st.form_submit_button("✅ Apply All Assignments", type="primary")
                    
                    if submitted:
                        rows_to_delete_mask = merged_df['Event Name'].isin([e for e, act in assignments.items() if act == '[DELETE]'])
                        
                        # Apply Deletions
                        if rows_to_delete_mask.sum() > 0:
                            merged_df = merged_df[~rows_to_delete_mask]
                            
                        # Apply Updates
                        for evt, cat in assignments.items():
                            if cat != '[DELETE]':
                                merged_df.loc[merged_df['Event Name'] == evt, 'Sub-Category'] = cat
                        
                        # Update Session State
                        st.session_state['data'] = merged_df
                        st.success("✅ Changes applied! Recalculating...")
                        st.rerun()


        else:
            st.success("🎉 All events are categorized!")



        # --- FINALIZE BUTTON ---
        st.markdown("---")
        st.markdown("### ✅ Ready to Analyze?")
        
        # Block progression if uncategorized items exist
        uncategorized_count = (merged_df['Sub-Category'] == 'uncategorized').sum()
        
        if uncategorized_count > 0:
            st.error(f"❌ **Cannot proceed**: {uncategorized_count} uncategorized records remain. Please use the Manual Resolution Tool above to assign categories or delete these records.")
            st.button("Finish Data Prep & Go to Analysis", type="primary", use_container_width=True, disabled=True)
        else:
            # THIS IS THE BUTTON USER WANTS TO AUTO NAVIGATE
            if st.button("Finish Data Prep & Go to Analysis", type="primary", use_container_width=True):
                # st.session_state['review_mode'] = False  <-- REMOVED to maintain state on Page 1 return
                st.session_state['current_page'] = 2 # Auto go to Page 2
                st.rerun()






# --- CONFIGURATION: EXCLUSION TERMS ---
EXCLUSION_TERMS = [
    "UOL", "UOB", "Warwick", "Stirling", "RMIT", "UOW", "La Trobe", "Sydney", "Monash", "UB", "Alberta", "GEM",
    "University of London", "University of Birmingham", "The University of Warwick", "University of Stirling",
    "RMIT University", "University of Wollongong", "La Trobe University", "The University of Sydney",
    "Monash College", "University at Buffalo, The State University of New York", "University of Alberta",
    "Grenoble Ecole de Management", "London", "Birmingham", "Wollongong", "Buffalo"
]

# --- PURPOSE DESCRIPTIONS FOR EACH QUESTION ---
QUESTION_PURPOSES = {
    1: "📈 **Analyzes specific attendance and attrition** trends across years to understand engagement and drop-off patterns",
    2: "👥 **Tracks student reach vs. retention** to distinguish between total volume and unique individual engagement",
    3: "🕒 **Studies optimal workshop timing** by analyzing attendance patterns across different days of the week and hours of the day",
    4: "🎓 **Compares university representation** to understand engagement of student in respective partner institutions",
    5: "📚 **Breaks down workshop categories** to identify which topics attract the most students",
    6: "🌍 **Analyzes local vs international** student participation to understand engagement differences between student types",
    7: "🔍 **Cross-analyzes workshop categories** by university to see which topics appeal to students from different institutions",
    8: "📊 **Cross-analyzes workshop categories** by academic major to understand which program types engage with which workshop topics",
    9: "⏱️ **Examines graduation proximity** to understand when students attend workshops relative to their expected graduation date",
    10: "📅 **Analyzes attrition rates** by comparing registered vs. actual attendance based on how far in advance students registered",
    11: "🏆 **Identifies top participating students** per university to recognize high engagement",
    12: "📅 **Tracks monthly attendance trends** by university to identify seasonal engagement patterns",
    13: "👨‍🏫 **Analyzes attendance by Trainer** to identify top-performing instructors",
    14: "🏷️ **Ranks workshops by attendance** to identify the most popular topics"
}

# --- FUNCTION TO RENDER SANDBOX BLOCKS ---
def render_sandbox(q_id, title, default_code, editable_title=False):
    st.markdown("---")
    
    # Create title row with Algorithm button on the right
    col_title, col_algo, col_run = st.columns([3, 1, 1])
    
    with col_title:
        if editable_title:
            # Use session state to persist custom titles
            key = f"title_q{q_id}"
            if key not in st.session_state:
                st.session_state[key] = title
            
            new_title = st.text_input(f"Q{q_id} Title", value=st.session_state[key], key=f"input_{key}")
            st.session_state[key] = new_title
            st.subheader(f"Q{q_id}: {new_title}")
        else:
            st.subheader(f"Q{q_id}: {title}")
    
    with col_algo:
        # Initialize session state for code visibility
        code_visibility_key = f"show_code_{q_id}"
        if code_visibility_key not in st.session_state:
            st.session_state[code_visibility_key] = False
        
        st.write("")  # Spacer for alignment
        # Algorithm button to toggle code visibility
        if st.button(f"🔍 Algorithm", key=f"algo_btn_{q_id}", help="View/Edit Python Logic", use_container_width=True):
            st.session_state[code_visibility_key] = not st.session_state[code_visibility_key]
    
    with col_run:
        st.write("")  # Spacer for alignment
        run_btn = st.button(f"▶ Run Q{q_id}", key=f"btn_{q_id}", type="primary", use_container_width=True)
    
    # Initialize session state for edited code
    edited_code_key = f"edited_code_{q_id}"
    if edited_code_key not in st.session_state:
        st.session_state[edited_code_key] = default_code
    
    # Toggle between Purpose view and Code Editor view
    if st.session_state[code_visibility_key]:
        # CODE EDITOR VIEW (Centered)
        st.markdown("")
        col_spacer1, col_center, col_spacer2 = st.columns([0.5, 4, 0.5])
        
        with col_center:
            st.markdown("#### 💻 Algorithm Editor")
            
            # Editable code text area
            edited_code = st.text_area(
                f"Python Logic for Q{q_id}", 
                value=st.session_state[edited_code_key], 
                height=350, 
                key=f"code_editor_{q_id}"
            )
            
            # Save and Reset buttons
            col_save, col_reset = st.columns([1, 1])
            
            with col_save:
                if st.button("💾 Save Changes", key=f"save_{q_id}", type="primary", use_container_width=True):
                    st.session_state[edited_code_key] = edited_code
                    st.success("✅ Code saved! Click Run to execute.", icon="✅")
            
            with col_reset:
                if st.button("🔄 Reset to Default", key=f"reset_{q_id}", use_container_width=True):
                    st.session_state[edited_code_key] = default_code
                    st.info("ℹ️ Code reset to default. Refresh to see changes.", icon="ℹ️")
                    st.rerun()
        
        # Use the saved edited code for execution
        code_input = st.session_state[edited_code_key]
    
    else:
        # PURPOSE VIEW (Default)
        purpose_text = QUESTION_PURPOSES.get(q_id, "Analyze workshop attendance data.")
        st.info(purpose_text)
        
        # Use default code (or previously saved edits) for execution
        code_input = st.session_state[edited_code_key]
    
    # Show attribute info and exclusion option in a collapsible section
    with st.expander("⚙️ Settings & Attributes", expanded=False):
        col_left, col_right = st.columns([3, 2])
        
        # Attribute Viewer (Specific & Compact)
        attr_map = {
            1: ['Attendance Status', 'SIMID', 'Event Name', 'Workshop Timing_Year'],
            2: ['Attendance Status', 'SIMID', 'Workshop Timing_Year'],
            3: ['Attendance Status', 'Workshop Timing_Day', 'Workshop Timing_Hours', 'Workshop Timing_Year'],
            4: ['Attendance Status', 'University Program', 'Workshop Timing_Year'],
            5: ['Attendance Status', 'Sub-Category', 'Workshop Timing_Year'],
            6: ['Attendance Status', 'Student_Type', 'Workshop Timing_Year'],
            7: ['Attendance Status', 'Sub-Category', 'University Program', 'Uni_Clean', 'Workshop Timing_Year'],
            8: ['Attendance Status', 'Sub-Category', 'Program Classification'],
            9: ['Attendance Status', 'Expected Grad Term', 'Grad_Year', 'Grad_Month', 'Grad_YYYY-MM', 'Attended Date'],
            10: ['Attendance Status', 'Registered Date', 'Workshop Timing_Year', 'Workshop Timing_Month', 'Workshop Timing_DayNumber'],
            11: ['Attendance Status', 'University Program', 'Student Name', 'SIMID', 'Display_Name'],
            12: ['Attendance Status', 'University Program', 'Uni_Clean', 'Uni_Grouped', 'Workshop Timing_Year', 'Workshop Timing_Month'],
            13: ['Attendance Status', 'Trainer'],
            14: ['Attendance Status', 'Event Name', 'Attended Date', 'Workshop Timing_Year']
        }
        
        if q_id in attr_map:
            display_cols = attr_map[q_id]
            msg = "Attributes used:"
            
            with col_left:
                st.markdown(f"**{msg}**")
                st.code(display_cols, language=None)
            
            with col_right:
                # Exclusion Checkbox
                exclude_uni = st.checkbox("Exclude Uni Events", key=f"exclude_{q_id}", help="Exclude workshops designed specifically for particular university")
            
            # --- PER-QUESTION DATA QUALITY ALERT ---
            # Check if any required attribute has missing values in the global dataset
            if st.session_state.get('data') is not None:
                df_global = st.session_state['data']
                missing_in_q = []
                for col in display_cols:
                    # Skip system-generated columns like ID
                    if col in ['ID']:
                        continue
                    
                    if col in df_global.columns:
                         n_missing = df_global[col].isnull().sum()
                         if n_missing > 0:
                             # Get the IDs of records with missing values
                             if 'ID' in df_global.columns:
                                 missing_ids = df_global[df_global[col].isnull()]['ID'].tolist()
                                 # Limit to first 10 IDs to avoid clutter
                                 if len(missing_ids) <= 10:
                                     ids_str = ", ".join(map(str, missing_ids))
                                 else:
                                     ids_str = ", ".join(map(str, missing_ids[:10])) + f"... (+{len(missing_ids)-10} more)"
                                 missing_in_q.append(f"{col} ({n_missing}) - IDs: {ids_str}")
                             else:
                                 missing_in_q.append(f"{col} ({n_missing})")
                    else:
                        missing_in_q.append(f"{col} (Missing)")
                
                if missing_in_q:
                    st.warning(f"⚠️ Potential missing data: {', '.join(missing_in_q)}")
        else:
            if st.session_state.get('data') is not None:
                display_cols = list(st.session_state['data'].columns)
            else:
                display_cols = []
            msg = "All attributes:"
            
            with col_left:
                st.markdown(f"**{msg}**")
                st.code(display_cols, language=None)
            
            with col_right:
                # Exclusion Checkbox
                exclude_uni = st.checkbox("Exclude Uni Events", key=f"exclude_{q_id}", help="Exclude workshops designed specifically for particular university")

    if run_btn:
        if st.session_state['data'] is None:
            st.error("Please upload data in Box A & B first.")
            return

        # Output Containers
        st.markdown("#### Results")
        
        # Prepare Data
        df_to_use = st.session_state['data'].copy()
        
        if exclude_uni:
            # Filter out events containing any exclusion term (case-insensitive)
            pattern = '|'.join([re.escape(term) for term in EXCLUSION_TERMS])
            # Ensure 'Event Name' exists
            if 'Event Name' in df_to_use.columns:
                mask = df_to_use['Event Name'].astype(str).str.contains(pattern, case=False, na=False)
                df_to_use = df_to_use[~mask]
                st.toast(f"ℹ️ Filtered out {mask.sum()} university-related events.", icon="🎓")
            else:
                st.warning("⚠️ 'Event Name' column missing. Cannot apply exclusion filter.")

        # Execution Environment
        local_vars = {
            'df': df_to_use, 
            'pd': pd, 
            'plt': plt, 
            'sns': sns,
            'np': __import__('numpy'),
            'kpi_result': {},
            'fig': None,
            'df_table': None,
            'figures_list': []
        }
        
        try:
            # Execute User Code
            exec(code_input, globals(), local_vars)
            
            # --- 1. Top Section: Years Compared ---
            results = local_vars.get('kpi_result', {})
            if isinstance(results, dict) and 'Years Compared' in results:
                st.markdown(f"""
                <div style="background-color: #e6f3ff; padding: 10px; border-radius: 5px; margin-bottom: 20px; border: 1px solid #b3d7ff;">
                    <strong style="color: #0056b3;">📅 Years Compared:</strong> <span style="color: #333;">{results['Years Compared']}</span>
                </div>
                """, unsafe_allow_html=True)
            elif isinstance(results, dict) and 'Status' in results:
                 st.info(results['Status'])

            # --- 2. Graphs & Tables Groups ---
            figures_list = local_vars.get('figures_list', [])
            fig = local_vars.get('fig')
            df_table = local_vars.get('df_table')
            
            # Helper to render a single group
            def render_group(table, figure, title=None):
                if title:
                    st.markdown(f"##### {title}")
                
                c1, c2 = st.columns([1, 1])
                
                # Left: Table
                with c1:
                    # st.markdown("#### 📝 Analysis Report")
                    if table is not None and not table.empty:
                        st.dataframe(table, use_container_width=True, hide_index=True)
                        
                        # Download table as high-res image
                        col_dl1, col_dl2 = st.columns([1, 1])
                        
                        try:
                            # 1. Image Download
                            # Create a matplotlib figure for the table
                            fig_table, ax_table = plt.subplots(figsize=(10, len(table) * 0.5 + 1))
                            ax_table.axis('tight')
                            ax_table.axis('off')
                            
                            # Create table
                            mpl_table = ax_table.table(
                                cellText=table.values,
                                colLabels=table.columns,
                                cellLoc='left',
                                loc='center'
                            )
                            
                            # Style the table
                            mpl_table.auto_set_font_size(False)
                            mpl_table.set_fontsize(10)
                            mpl_table.scale(1, 2)
                            
                            # Header styling
                            for i in range(len(table.columns)):
                                cell = mpl_table[(0, i)]
                                cell.set_facecolor('#4472C4')
                                cell.set_text_props(weight='bold', color='white')
                            
                            plt.tight_layout()
                            
                            # Save to buffer
                            buf_table = io.BytesIO()
                            fig_table.savefig(buf_table, format='png', dpi=300, bbox_inches='tight', facecolor='white')
                            buf_table.seek(0)
                            plt.close(fig_table)
                            
                            with col_dl1:
                                st.download_button(
                                    label="📥 Download Image",
                                    data=buf_table,
                                    file_name=f"table_{q_id}_{int(pd.Timestamp.now().timestamp())}.png",
                                    mime="image/png",
                                    key=f"dl_table_img_{q_id}_{int(pd.Timestamp.now().timestamp())}_{np.random.randint(0,10000)}"
                                )
                                
                            # 2. Excel Download (Auto-width)
                            buf_xlsx = io.BytesIO()
                            with pd.ExcelWriter(buf_xlsx, engine='xlsxwriter') as writer:
                                table.to_excel(writer, sheet_name='Data', index=False)
                                worksheet = writer.sheets['Data']
                                for i, col in enumerate(table.columns):
                                    # Calculate width: max(header_len, max_cell_len) + padding
                                    max_len = len(str(col))
                                    column_data = table[col].astype(str)
                                    if not column_data.empty:
                                        max_len = max(max_len, column_data.map(len).max())
                                    worksheet.set_column(i, i, max_len + 2)
                            
                            buf_xlsx.seek(0)
                            
                            with col_dl2:
                                st.download_button(
                                    label="📥 Download Excel",
                                    data=buf_xlsx,
                                    file_name=f"table_{q_id}_{int(pd.Timestamp.now().timestamp())}.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                    key=f"dl_table_xlsx_{q_id}_{int(pd.Timestamp.now().timestamp())}_{np.random.randint(0,10000)}"
                                )

                        except Exception as e:
                            st.caption(f"Download unavailable: {e}")
                    else:
                        st.caption("No table data available.")
                
                # Right: Graph
                with c2:
                    if figure:
                        st.pyplot(figure)
                        # Download button
                        buf = io.BytesIO()
                        figure.savefig(buf, format='png', dpi=300, bbox_inches='tight')
                        buf.seek(0)
                        st.download_button(
                            label="📥 Download Graph",
                            data=buf,
                            file_name=f"graph_{q_id}_{int(pd.Timestamp.now().timestamp())}.png",
                            mime="image/png",
                            key=f"dl_{q_id}_{int(pd.Timestamp.now().timestamp())}_{np.random.randint(0,1000)}"
                        )

            # Case A: Multiple Figures (from figures_list)
            if figures_list:
                for i, item in enumerate(figures_list):
                    f = item.get('fig')
                    t = item.get('table') # Expecting table in the item now
                    title = item.get('title')
                    
                    # If table not explicitly in item, maybe use global df_table for the first one?
                    # For now, let's assume Q3/Q5 updates will put it there.
                    # Fallback: if no table in item, use empty
                    
                    render_group(t, f, title)
                    st.markdown("---")

            # Case B: Single Figure (Backward Compatibility)
            elif fig or (df_table is not None and not df_table.empty):
                render_group(df_table, fig)

                    
        except Exception as e:
            st.error(f"❌ Logic Error: {e}")
    else:
        st.info("Click 'Run' to execute analysis.")

# ==============================================================================
# DEFAULT CODE BLOCKS (Pre-filled Logic)
# ==============================================================================

# Q1 CODE
# Q1 CODE
# Q1 CODE
# Q1 CODE
code_q1 = """
# 1. Setup
figures_list = []
kpi_result = {}

# 2. Logic: Prepare Data
# Create working copy
df_calc = df.copy()

# Clean Year
if 'Date' in df_calc.columns:
    df_calc['Workshop Timing_Year'] = pd.to_datetime(df_calc['Date'], errors='coerce').dt.year

df_calc = df_calc.dropna(subset=['Workshop Timing_Year'])
df_calc['Workshop Timing_Year'] = df_calc['Workshop Timing_Year'].astype(int)

# Identify Attended
df_calc['Is_Attended'] = df_calc['Attendance Status'].astype(str).str.lower() == 'attended'

# Group by Year
yearly_metrics = df_calc.groupby('Workshop Timing_Year').agg(
    Registered=('ID', 'size'),
    Attended=('Is_Attended', 'sum')
)

# Calculate Attrition Count
yearly_metrics['Attrition Count'] = yearly_metrics['Registered'] - yearly_metrics['Attended']
years = yearly_metrics.index.sort_values()

# Set KPI Result with Years Compared
if not years.empty:
    kpi_result["Years Compared"] = f"{years.min()} - {years.max()}"


# --- SECTION 2: COMBINED TRENDS GRAPH ---
if not yearly_metrics.empty:
    fig_comb, ax_comb = plt.subplots(figsize=(10, 6))
    
    # Plot both lines
    ax_comb.plot(years, yearly_metrics['Attended'], marker='o', linewidth=2, markersize=8, color='#2ca02c', label='Attended')
    ax_comb.plot(years, yearly_metrics['Attrition Count'], marker='o', linewidth=2, markersize=8, color='#d62728', label='Attrition Count')
    
    ax_comb.set_title('Attendance vs Attrition Trends', fontsize=14, weight='bold')
    ax_comb.set_xlabel('Year', fontsize=12)
    ax_comb.set_ylabel('Count', fontsize=12)
    ax_comb.grid(True, alpha=0.3, linestyle='--')
    ax_comb.xaxis.set_major_locator(plt.MaxNLocator(integer=True))
    # Extend Y-axis by 15% to prevent overlap
    if not yearly_metrics[['Attended', 'Attrition Count']].empty:
        max_val = yearly_metrics[['Attended', 'Attrition Count']].max().max()
        ax_comb.set_ylim(0, max_val * 1.15)
    ax_comb.legend()
    
    # Add labels for both
    for year, count in zip(years, yearly_metrics['Attended']):
        ax_comb.annotate(str(count), (year, count), textcoords="offset points", xytext=(0,10), ha='center', fontsize=9, weight='bold', color='#2ca02c')
        
    for year, count in zip(years, yearly_metrics['Attrition Count']):
        ax_comb.annotate(str(count), (year, count), textcoords="offset points", xytext=(0,-15), ha='center', fontsize=9, weight='bold', color='#d62728')

    # Create Combined Statistics Table
    df_combined_table = yearly_metrics[['Registered', 'Attended', 'Attrition Count']].reset_index()
    df_combined_table.columns = ['Year', 'Registered', 'Attended', 'Attrition Count']
    
    # Add % Change columns
    if len(df_combined_table) >= 2:
        for i in range(len(df_combined_table) - 1):
            y1, y2 = df_combined_table.iloc[i]['Year'], df_combined_table.iloc[i+1]['Year']
            
            # Registered % Change
            r1 = df_combined_table.iloc[i]['Registered']
            r2 = df_combined_table.iloc[i+1]['Registered']
            reg_pct = f"{((r2 - r1) / r1) * 100:+.1f}%" if r1 != 0 else "New"
            df_combined_table.loc[i, f'% Chg Reg ({int(y1)}->{int(y2)})'] = reg_pct
            
            # Attended % Change
            a1 = df_combined_table.iloc[i]['Attended']
            a2 = df_combined_table.iloc[i+1]['Attended']
            att_pct = f"{((a2 - a1) / a1) * 100:+.1f}%" if a1 != 0 else "New"
            df_combined_table.loc[i, f'% Chg Att ({int(y1)}->{int(y2)})'] = att_pct
            
            # Attrition % Change
            at1 = df_combined_table.iloc[i]['Attrition Count']
            at2 = df_combined_table.iloc[i+1]['Attrition Count']
            atr_pct = f"{((at2 - at1) / at1) * 100:+.1f}%" if at1 != 0 else "New"
            df_combined_table.loc[i, f'% Chg Atr ({int(y1)}->{int(y2)})'] = atr_pct
        
        # Fill NaN for last row
        df_combined_table = df_combined_table.fillna('-')

    figures_list.append({
        'title': f'Combined Trends ({years.min()}-{years.max()})',
        'table': df_combined_table, 
        'fig': fig_comb
    })


# --- SECTION 3: ATTENDANCE ANALYSIS (Line Graph + Table) ---
if not yearly_metrics.empty:
    # 2a. Table
    df_att = yearly_metrics[['Attended']].copy().reset_index()
    df_att.rename(columns={'Workshop Timing_Year': 'Year'}, inplace=True)
    
    # Calculate % Change
    df_att['% Change'] = df_att['Attended'].pct_change().mul(100).fillna(0)
    # Format % Change (handle first row as New or -)
    def fmt_pct(val, idx, series):
        if idx == 0: return "-"
        prev = series.iloc[idx-1]
        if prev == 0: return "New"
        return f"{val:+.1f}%"
        
    df_att['% Change'] = [fmt_pct(x, i, df_att['Attended']) for i, x in enumerate(df_att['% Change'])]
    
    # 2b. Graph
    fig_att, ax_att = plt.subplots(figsize=(10, 6))
    ax_att.plot(years, yearly_metrics['Attended'], marker='o', linewidth=2, markersize=8, color='#2ca02c', label='Attended')
    
    ax_att.set_title('Yearly Attendance Trend', fontsize=14, weight='bold')
    ax_att.set_xlabel('Year', fontsize=12)
    ax_att.set_ylabel('Attendees', fontsize=12)
    ax_att.grid(True, alpha=0.3, linestyle='--')
    ax_att.xaxis.set_major_locator(plt.MaxNLocator(integer=True))
    
    # Extend Y-axis by 15%
    if not yearly_metrics['Attended'].empty:
        ax_att.set_ylim(0, yearly_metrics['Attended'].max() * 1.15)
    
    # Add labels
    for year, count in zip(years, yearly_metrics['Attended']):
        ax_att.annotate(str(count), (year, count), textcoords="offset points", xytext=(0,10), ha='center', fontsize=9, weight='bold')
        
    figures_list.append({
        'title': f'Yearly Attendance Analysis ({years.min()}-{years.max()})',
        'table': df_att,
        'fig': fig_att
    })


# --- SECTION 3: ATTRITION ANALYSIS (Line Graph + Table) ---
if not yearly_metrics.empty:
    # 3a. Table
    df_attr = yearly_metrics[['Attrition Count']].copy().reset_index()
    df_attr.rename(columns={'Workshop Timing_Year': 'Year'}, inplace=True)
    
    # Calculate % Change
    df_attr['% Change'] = df_attr['Attrition Count'].pct_change().mul(100).fillna(0)
    
    df_attr['% Change'] = [fmt_pct(x, i, df_attr['Attrition Count']) for i, x in enumerate(df_attr['% Change'])]
    
    # 3b. Graph
    fig_attr, ax_attr = plt.subplots(figsize=(10, 6))
    ax_attr.plot(years, yearly_metrics['Attrition Count'], marker='o', linewidth=2, markersize=8, color='#d62728', label='Attrition Count')
    
    ax_attr.set_title('Yearly Attrition Trend', fontsize=14, weight='bold')
    ax_attr.set_xlabel('Year', fontsize=12)
    ax_attr.set_ylabel('Attrition Count', fontsize=12)
    ax_attr.grid(True, alpha=0.3, linestyle='--')
    ax_attr.xaxis.set_major_locator(plt.MaxNLocator(integer=True))
    
    # Extend Y-axis by 15%
    if not yearly_metrics['Attrition Count'].empty:
        ax_attr.set_ylim(0, yearly_metrics['Attrition Count'].max() * 1.15)
    
    # Add labels
    for year, count in zip(years, yearly_metrics['Attrition Count']):
        ax_attr.annotate(str(count), (year, count), textcoords="offset points", xytext=(0,10), ha='center', fontsize=9, weight='bold')

    figures_list.append({
        'title': f'Yearly Attrition Analysis ({years.min()}-{years.max()})',
        'table': df_attr,
        'fig': fig_attr
    })
else:
    kpi_result['Status'] = "No valid yearly data found."
"""


# Q2 CODE
code_q2 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# 2. Logic: Unique count per year
unique_per_year = df_attended.groupby('Workshop Timing_Year')['SIMID'].nunique()

# 3. Stats & Table Construction
total_unique = df_attended['SIMID'].nunique()
years = sorted(unique_per_year.index)

# Build Row Data
row_data = {
    "Metric": "Unique Students",
    "Total": total_unique
}

# Add Yearly Counts
for year in years:
    row_data[str(year)] = unique_per_year[year]

# Add % Changes
if len(years) >= 2:
    for i in range(len(years) - 1):
        y1, y2 = years[i], years[i+1]
        c1 = unique_per_year[y1]
        c2 = unique_per_year[y2]
        
        if c1 == 0:
            pct = "New" if c2 > 0 else "-"
        else:
            pct = f"{((c2 - c1) / c1) * 100:+.1f}%"
        
        row_data[f"% Change ({y1}->{y2})"] = pct

# Create DataFrame
df_table = pd.DataFrame([row_data])

# Set KPI Result with Years Compared
kpi_result = {
    "Years Compared": ", ".join(map(str, years))
}

# 4. Graph
fig, ax = plt.subplots(figsize=(10, 6))
ax.plot(unique_per_year.index, unique_per_year.values, marker='o', linewidth=2, markersize=8, color='#ff7f0e')
ax.set_title('Unique Participant Count per Year', fontsize=14, weight='bold')
ax.set_ylabel('Number of Unique Students', fontsize=11)
ax.set_xlabel('Year', fontsize=11)
ax.grid(True, alpha=0.3, linestyle='--')

# Format x-axis to show only integer years
ax.xaxis.set_major_locator(plt.MaxNLocator(integer=True))

# Extend Y-axis by 15%
if not unique_per_year.empty:
    ax.set_ylim(0, unique_per_year.max() * 1.15)

# Add labels
for year, count in unique_per_year.items():
    ax.text(year, count, str(count), ha='center', va='bottom', fontsize=9)

plt.tight_layout()

# Clear bottom report
text_report = ""
"""

# Q3 CODE
code_q3 = """
# 1. Filter for Attended records
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# Initialize defaults
kpi_result = {}
figures_list = []
fig = None
df_table = pd.DataFrame()

# Normalize and map day names
def normalize_day(day_str):
    # Define mappings inside function to avoid scope issues when executed as string
    days_full = {
        'monday': 'Monday',
        'tuesday': 'Tuesday', 
        'wednesday': 'Wednesday',
        'thursday': 'Thursday',
        'friday': 'Friday',
        'saturday': 'Saturday',
        'sunday': 'Sunday'
    }
    
    day_variations = {
        # Monday
        'mon': 'monday', 'mond': 'monday',
        # Tuesday
        'tue': 'tuesday', 'tues': 'tuesday', 'tu': 'tuesday',
        # Wednesday
        'wed': 'wednesday', 'weds': 'wednesday',
        # Thursday
        'thu': 'thursday', 'thur': 'thursday', 'thurs': 'thursday',
        # Friday
        'fri': 'friday',
        # Saturday
        'sat': 'saturday',
        # Sunday
        'sun': 'sunday'
    }
    
    if pd.isna(day_str):
        return None
    
    # Clean: lowercase, strip whitespace
    clean = str(day_str).lower().strip()
    
    # Direct match with full name
    if clean in days_full:
        return days_full[clean]
    
    # Match with abbreviations
    if clean in day_variations:
        return days_full[day_variations[clean]]
    
    # Fuzzy match (for typos)
    import difflib
    all_day_keys = list(days_full.keys()) + list(day_variations.keys())
    matches = difflib.get_close_matches(clean, all_day_keys, n=1, cutoff=0.7)
    
    if matches:
        matched_key = matches[0]
        if matched_key in days_full:
            return days_full[matched_key]
        elif matched_key in day_variations:
            return days_full[day_variations[matched_key]]
    
    return None

# 2. Validate required columns exist
if 'Workshop Timing_Day' not in df_attended.columns or 'Workshop Timing_Hours' not in df_attended.columns:
    missing = []
    if 'Workshop Timing_Day' not in df_attended.columns: missing.append('Workshop Timing_Day')
    if 'Workshop Timing_Hours' not in df_attended.columns: missing.append('Workshop Timing_Hours')
    
    kpi_result['Status'] = f"Missing required columns: {', '.join(missing)}"
    df_table = pd.DataFrame([{"Error": kpi_result['Status']}])
else:
    # Filter for valid timing data
    df_attended = df_attended.dropna(subset=['Workshop Timing_Day', 'Workshop Timing_Hours'])
    
    if df_attended.empty:
        kpi_result['Status'] = "No valid timing data found for attended events."
        df_table = pd.DataFrame([{"Error": "No valid timing data", "Debug": f"Attended records: {len(df[df['Attendance Status'].astype(str).str.lower() == 'attended'])}"}])
    else:
        # Convert Hour to integer
        df_attended['Hour'] = pd.to_numeric(df_attended['Workshop Timing_Hours'], errors='coerce').astype('Int64')
        
        # Check if hours are valid
        valid_hours = df_attended['Hour'].notna()
        if not valid_hours.any():
            kpi_result['Status'] = "No valid hour data found."
            df_table = pd.DataFrame([{"Error": "Invalid hour format", "Sample Hours": str(df_attended['Workshop Timing_Hours'].head(3).tolist())}])
        else:
            df_attended = df_attended[valid_hours]
            
            # Apply normalization
            df_attended['Day_Normalized'] = df_attended['Workshop Timing_Day'].apply(normalize_day)
            
            # Filter out rows where day couldn't be normalized
            df_attended = df_attended[df_attended['Day_Normalized'].notna()]
            
            if df_attended.empty:
                kpi_result['Status'] = "No valid day names found."
                df_table = pd.DataFrame([{"Error": "Invalid day names", "Sample Days": str(df[df['Attendance Status'].astype(str).str.lower() == 'attended']['Workshop Timing_Day'].unique()[:5].tolist())}])
            else:
                # Create ordered categorical with proper case
                days_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
                df_attended['Day_Ordered'] = pd.Categorical(
                    df_attended['Day_Normalized'], 
                    categories=days_order, 
                    ordered=True
                )
                
                # Update Workshop Timing_Day to use normalized values
                df_attended['Workshop Timing_Day'] = df_attended['Day_Normalized']
        
                # 3. Overall Analysis
                # Top 3 Time Slots
                overall_slots = df_attended.groupby(['Workshop Timing_Day', 'Hour']).size().sort_values(ascending=False).head(3)
                
                overall_data = []
                for rank, ((day, hour), count) in enumerate(overall_slots.items(), 1):
                    # Drill down: Get Top 10 Workshops for this slot
                    slot_mask = (df_attended['Workshop Timing_Day'] == day) & (df_attended['Hour'] == hour)
                    df_slot = df_attended[slot_mask]
                    
                    # Group by Event Name
                    # Count Attendance (size) and Runs (unique Attended Date)
                    ws_stats = df_slot.groupby('Event Name').agg(
                        Attendance=('SIMID', 'size'),
                        Runs=('Attended Date', 'nunique')
                    ).sort_values('Attendance', ascending=False).head(10)
                    
                    # Format details
                    details_list = []
                    for idx, (event, row) in enumerate(ws_stats.iterrows(), 1):
                        run_str = f", {row['Runs']} runs" if row['Runs'] > 1 else ", 1 run"
                        details_list.append(f"{idx}. {event} ({row['Attendance']} pax{run_str})")
                    
                    # Join with newlines
                    details_str = "\\n".join(details_list)

                    overall_data.append({
                        'Rank': rank,
                        'Day': day,
                        'Time Slot': f"{int(hour):02d}:00",
                        'Attendance Count': count,
                        'Top 10 Workshops': details_str
                    })
                
                df_overall_table = pd.DataFrame(overall_data)
                
                # Overall Heatmap (Previous logic)
                pivot_overall = df_attended.pivot_table(
                    index='Day_Ordered', 
                    columns='Hour', 
                    values='SIMID', 
                    aggfunc='count', 
                    fill_value=0
                )
                
                if not pivot_overall.empty:
                    try:
                        fig, ax = plt.subplots(figsize=(14, 6))
                        sns.heatmap(pivot_overall, cmap="YlGnBu", annot=True, fmt="d", ax=ax, cbar_kws={'label': 'Attendance Count'})
                        ax.set_title('Overall: Workshop Attendance by Day and Hour', fontsize=14, weight='bold')
                        ax.set_xlabel('Hour of Day')
                        ax.set_ylabel('Day of Week')
                        plt.tight_layout()
                        
                        figures_list.append({
                            'fig': fig,
                            'title': 'Overall Heatmap',
                            'table': df_overall_table
                        })
                    except Exception as e:
                        pass
                
                # 4. Yearly Analysis
                years = sorted(df_attended['Workshop Timing_Year'].unique())
                
                yearly_data = []
                for year in years:
                    df_year = df_attended[df_attended['Workshop Timing_Year'] == year]
                    
                    if not df_year.empty:
                        # Top 3 Time Slots for this year
                        year_slots = df_year.groupby(['Workshop Timing_Day', 'Hour']).size().sort_values(ascending=False).head(3)
                        
                        for rank, ((day, hour), count) in enumerate(year_slots.items(), 1):
                            # Drill down for Year
                            slot_mask = (df_year['Workshop Timing_Day'] == day) & (df_year['Hour'] == hour)
                            df_slot = df_year[slot_mask]
                            
                            ws_stats = df_slot.groupby('Event Name').agg(
                                Attendance=('SIMID', 'size'),
                                Runs=('Attended Date', 'nunique')
                            ).sort_values('Attendance', ascending=False).head(10)
                            
                            details_list = []
                            for idx, (event, row) in enumerate(ws_stats.iterrows(), 1):
                                run_str = f", {row['Runs']} runs" if row['Runs'] > 1 else ", 1 run"
                                details_list.append(f"{idx}. {event} ({row['Attendance']} pax{run_str})")
                            
                            details_str = "\\n".join(details_list)

                            yearly_data.append({
                                'Year': int(year),
                                'Rank': rank,
                                'Day': day,
                                'Time Slot': f"{int(hour):02d}:00",
                                'Attendance Count': count,
                                'Top 10 Workshops': details_str
                            })
                        
                        # Year-specific table
                        # For yearly specific table in figures_list, we reuse the same structure
                        df_year_table = pd.DataFrame(yearly_data[-len(year_slots):]) # Get last N added

                        # Yearly Heatmap
                        pivot_year = df_year.pivot_table(
                            index='Day_Ordered',
                            columns='Hour',
                            values='SIMID',
                            aggfunc='count',
                            fill_value=0
                        )
                        
                        if not pivot_year.empty:
                            try:
                                fig, ax = plt.subplots(figsize=(14, 6))
                                sns.heatmap(pivot_year, cmap="YlGnBu", annot=True, fmt="d", ax=ax, cbar_kws={'label': 'Attendance Count'})
                                ax.set_title(f'{int(year)}: Workshop Attendance by Day and Hour', fontsize=14, weight='bold')
                                ax.set_xlabel('Hour of Day')
                                ax.set_ylabel('Day of Week')
                                plt.tight_layout()
                                
                                figures_list.append({
                                    'fig': fig,
                                    'title': f'{int(year)} Heatmap',
                                    'table': df_year_table
                                })
                            except Exception as e:
                                pass
                
                # 5. Create combined summary table
                df_yearly_all = pd.DataFrame(yearly_data)
                
                # Add separator
                separator = pd.DataFrame([{'Year': '---', 'Rank': '---', 'Day': '---', 'Time Slot': '---', 'Attendance Count': '---', 'Top 10 Workshops': '---'}])
                
                # Overall section with Year column
                df_overall_display = df_overall_table.copy()
                df_overall_display.insert(0, 'Year', 'Overall')
                
                # Combine: Overall + Separator + Yearly
                df_table = pd.concat([df_overall_display, separator, df_yearly_all], ignore_index=True)
                
                # Set KPI
                if overall_data:
                    top_slot = overall_data[0]
                    kpi_result['Most Popular Slot'] = f"{top_slot['Day']} {top_slot['Time Slot']} ({top_slot['Attendance Count']} attendees)"
                else:
                    kpi_result['Status'] = 'No data available'
"""



# Q4 CODE
code_q4 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# 2. Logic
# Clean University Name (take first part before comma if applicable) and Title Case
df_attended['Uni_Clean'] = df_attended['University Program'].astype(str).apply(lambda x: x.split(',')[0]).str.title()

# Identify Top 10 Universities Overall for Graph
top_10_unis = df_attended['Uni_Clean'].value_counts().head(10).index
df_top_10 = df_attended[df_attended['Uni_Clean'].isin(top_10_unis)].copy()

# Aggregate by University and Year for Graph
uni_year_counts = df_top_10.groupby(['Uni_Clean', 'Workshop Timing_Year']).size().reset_index(name='Count')

# 3. Stats
kpi_result = {
    "Years Compared": ", ".join(map(str, sorted(df_attended['Workshop Timing_Year'].unique())))
}

# 4. Graph
fig, ax = plt.subplots(figsize=(12, 8))
sns.lineplot(
    data=uni_year_counts, 
    x='Workshop Timing_Year', 
    y='Count', 
    hue='Uni_Clean', 
    marker='o',
    palette='viridis', 
    hue_order=top_10_unis, # Maintain overall top 10 order
    ax=ax
)
ax.set_title('Top 10 Universities: Attendance Trends by Year')
ax.set_xlabel('Year')
ax.set_ylabel('Number of Attendees')
ax.grid(True, alpha=0.3, linestyle='--')
ax.xaxis.set_major_locator(plt.MaxNLocator(integer=True))

# Extend Y-axis by 15%
if not uni_year_counts.empty:
     ax.set_ylim(0, uni_year_counts['Count'].max() * 1.15)
     
ax.legend(title='University', bbox_to_anchor=(1.02, 1), loc='upper left')
plt.tight_layout()

# 5. Table & Description
# Aggregate ALL universities
all_uni_counts = df_attended.groupby(['Uni_Clean', 'Workshop Timing_Year']).size().unstack(fill_value=0)
# Sort by total
all_uni_counts['Total'] = all_uni_counts.sum(axis=1)
all_uni_counts = all_uni_counts.sort_values('Total', ascending=False).drop(columns='Total')

# Calculate % Increase
years = sorted(df_attended['Workshop Timing_Year'].unique())
if len(years) >= 2:
    for i in range(len(years) - 1):
        y1, y2 = years[i], years[i+1]
        col_name = f'% Change ({y1}->{y2})'
        
        all_uni_counts[col_name] = all_uni_counts.apply(
            lambda row, start=y1, end=y2: (
                f"{((row[end] - row[start]) / row[start]) * 100:+.1f}%" 
                if row[start] != 0 
                else ("New" if row[end] > 0 else "-")
            ), 
            axis=1
        )

# Reset index to make University a column
df_table = all_uni_counts.reset_index()

text_report = ""
"""

# Q5 CODE
# Q5 CODE
# Q5 CODE
code_q5 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# 2. Logic: Attendance by Sub-Category
sub_cat_counts = df_attended['Sub-Category'].value_counts()

# 3. Stats
kpi_result = {
    "Years Compared": ", ".join(map(str, sorted(df_attended['Workshop Timing_Year'].unique())))
}

figures_list = []
df_table = pd.DataFrame()

# 4. Graph: Overall Pie Chart
if not sub_cat_counts.empty:
    try:
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Create pie chart
        # Use pastel colormap
        cmap = plt.get_cmap('Pastel1')
        colors = [cmap(i) for i in range(len(sub_cat_counts))]
        
        # Custom autopct function to hide % when < 1%
        def autopct_format(pct):
            return f'{pct:.1f}%' if pct >= 1 else ''
        
        wedges, texts, autotexts = ax.pie(
            sub_cat_counts.values,
            labels=None,  # No labels on pie slices, use legend only
            autopct=autopct_format,
            startangle=90,
            colors=colors,
            pctdistance=0.85,
            explode=[0.05 if i == 0 else 0 for i in range(len(sub_cat_counts))]
        )
        
        
        # Draw circle for donut style
        centre_circle = plt.Circle((0,0),0.70,fc='white')
        fig.gca().add_artist(centre_circle)
        
        ax.set_title('Overall Workshop Attendance by Sub-Category', fontsize=14, weight='bold')
        ax.legend(
            [s.title() if isinstance(s, str) else s for s in sub_cat_counts.index],
            title="Sub-Category",
            loc="center left",
            bbox_to_anchor=(1, 0, 0.5, 1)
        )
        ax.axis('equal')
        plt.tight_layout()
        
        # Add Percentage to Table
        df_overall_table = sub_cat_counts.reset_index(name='Count').rename(columns={'index': 'Sub-Category'})
        df_overall_table['Percentage'] = (df_overall_table['Count'] / df_overall_table['Count'].sum() * 100).map('{:.1f}%'.format)
        # Apply Title Case
        df_overall_table['Sub-Category'] = df_overall_table['Sub-Category'].astype(str).str.title()
        
        figures_list.append({
            'fig': fig,
            'title': 'Overall Distribution',
            'table': df_overall_table
        })
        
    except Exception as e:
        kpi_result["Status"] = f"Error plotting Overall Pie Chart: {e}"
else:
    kpi_result["Status"] = "No valid Sub-Category data found."

# 5. Graph: Yearly Pie Charts
years = sorted(df_attended['Workshop Timing_Year'].unique())
for year in years:
    df_year = df_attended[df_attended['Workshop Timing_Year'] == year]
    if not df_year.empty:
        year_counts = df_year['Sub-Category'].value_counts()
        
        if not year_counts.empty:
            try:
                fig, ax = plt.subplots(figsize=(10, 6))
                
                # Create pie chart for year
                cmap = plt.get_cmap('Pastel1')
                colors = [cmap(i) for i in range(len(year_counts))]
                
                # Custom autopct function to hide % when < 1%
                def autopct_format(pct):
                    return f'{pct:.1f}%' if pct >= 1 else ''
                
                wedges, texts, autotexts = ax.pie(
                    year_counts.values,
                    labels=None,  # No labels on pie slices, use legend only
                    autopct=autopct_format,
                    startangle=90,
                    colors=colors,
                    pctdistance=0.85,
                    explode=[0.05 if i == 0 else 0 for i in range(len(year_counts))]
                )
                
                
                centre_circle = plt.Circle((0,0),0.70,fc='white')
                fig.gca().add_artist(centre_circle)
                
                ax.set_title(f'{year} Workshop Attendance by Sub-Category', fontsize=14, weight='bold')
                ax.legend(
                    [s.title() if isinstance(s, str) else s for s in year_counts.index],
                    title="Sub-Category",
                    loc="center left",
                    bbox_to_anchor=(1, 0, 0.5, 1)
                )
                ax.axis('equal')
                plt.tight_layout()
                
                # Add Percentage to Table
                df_year_table = year_counts.reset_index(name='Count').rename(columns={'index': 'Sub-Category'})
                df_year_table['Percentage'] = (df_year_table['Count'] / df_year_table['Count'].sum() * 100).map('{:.1f}%'.format)
                # Apply Title Case
                df_year_table['Sub-Category'] = df_year_table['Sub-Category'].astype(str).str.title()
                
                figures_list.append({
                    'fig': fig,
                    'title': f'{year} Distribution',
                    'table': df_year_table
                })
            except Exception as e:
                pass

# 6. Table: Overall Comparison
# Pivot to show counts per year
pivot_counts = df_attended.pivot_table(index='Sub-Category', columns='Workshop Timing_Year', values='SIMID', aggfunc='count', fill_value=0)
pivot_counts['Total'] = pivot_counts.sum(axis=1)

# Sort by Total
pivot_counts = pivot_counts.sort_values('Total', ascending=False)

# Add Overall Percentage
pivot_counts['Overall %'] = (pivot_counts['Total'] / pivot_counts['Total'].sum() * 100).map('{:.1f}%'.format)

# Calculate % Increase
years = sorted(df_attended['Workshop Timing_Year'].unique())
if len(years) >= 2:
    for i in range(len(years) - 1):
        y1, y2 = years[i], years[i+1]
        col_name = f'% Change ({y1}->{y2})'
        
        # Use default args to capture y1, y2 values to avoid scope issues in exec()
        pivot_counts[col_name] = pivot_counts.apply(
            lambda row, start=y1, end=y2: (
                f"{((row[end] - row[start]) / row[start]) * 100:+.1f}%" 
                if row[start] != 0 
                else ("New" if row[end] > 0 else "-")
            ), 
            axis=1
        )

df_table = pivot_counts.reset_index()
# Apply Title Case
df_table['Sub-Category'] = df_table['Sub-Category'].astype(str).str.title()
"""


# Q6 CODE
code_q6 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# 2. Logic
# Use existing Student_Type column if available, otherwise create it
if 'Student_Type' not in df_attended.columns:
    # Student Type: Use Citizenship, treat NaN as Local
    def get_student_type(val):
        # If NaN or empty, treat as Local
        if pd.isna(val) or str(val).strip() == '':
            return 'Local'
        val_str = str(val).lower()
        if 'singapore' in val_str or 'pr' in val_str:
            return 'Local'
        return 'International'

    if 'Citizenship' in df_attended.columns:
        df_attended['Student_Type'] = df_attended['Citizenship'].apply(get_student_type)
    else:
        df_attended['Student_Type'] = 'Local'

    # Force standardization to Title Case (Local, International) to match color map keys
    df_attended['Student_Type'] = df_attended['Student_Type'].astype(str).str.strip().str.title()

# Debug: Show Student_Type distribution
student_type_counts = df_attended['Student_Type'].value_counts()
kpi_result = {
    'Total Records': len(df_attended),
    'Local Count': student_type_counts.get('Local', 0),
    'International Count': student_type_counts.get('International', 0)
}

# 3. Stats
figures_list = []

# 4. Yearly Pie Charts
years = sorted(df_attended['Workshop Timing_Year'].unique())

for year in years:
    df_year = df_attended[df_attended['Workshop Timing_Year'] == year]
    
    if not df_year.empty:
        year_type_counts = df_year['Student_Type'].value_counts()
        
        if not year_type_counts.empty:
            try:
                fig, ax = plt.subplots(figsize=(10, 6))
                
                # Color mapping: Local = Blue, International = Pink
                colors_map = {'Local': '#1f77b4', 'International': '#e377c2'}
                colors = [colors_map.get(label, '#cccccc') for label in year_type_counts.index]
                
                # Custom autopct function to hide % when < 1%
                def autopct_format(pct):
                    return f'{pct:.1f}%' if pct >= 1 else ''
                
                wedges, texts, autotexts = ax.pie(
                    year_type_counts.values,
                    labels=None,  # No labels on pie slices, use legend only
                    autopct=autopct_format,
                    startangle=90,
                    colors=colors,
                    pctdistance=0.85,
                    explode=[0.05 if i == 0 else 0 for i in range(len(year_type_counts))]
                )
                
                # Draw circle for donut style
                centre_circle = plt.Circle((0,0),0.70,fc='white')
                fig.gca().add_artist(centre_circle)
                
                ax.set_title(f'{int(year)}: Local vs International Student Attendance', fontsize=14, weight='bold')
                ax.legend(
                    year_type_counts.index,
                    title="Student Type",
                    loc="center left",
                    bbox_to_anchor=(1, 0, 0.5, 1)
                )
                ax.axis('equal')
                plt.tight_layout()
                
                # Create table for this year
                df_year_table = year_type_counts.reset_index(name='Count').rename(columns={'index': 'Student Type'})
                df_year_table['Percentage'] = (df_year_table['Count'] / df_year_table['Count'].sum() * 100).map('{:.1f}%'.format)
                
                figures_list.append({
                    'fig': fig,
                    'title': f'{int(year)} Distribution',
                    'table': df_year_table
                })
            except Exception as e:
                pass

# 5. Overall Bar Graph (Yearly Comparison)
# Create a grouped bar chart showing Local vs International for each year
years = sorted(df_attended['Workshop Timing_Year'].unique())

# Prepare data for bar chart
bar_data = []
for year in years:
    df_year = df_attended[df_attended['Workshop Timing_Year'] == year]
    year_type_counts = df_year['Student_Type'].value_counts()
    
    local_count = year_type_counts.get('Local', 0)
    intl_count = year_type_counts.get('International', 0)
    
    bar_data.append({
        'Year': int(year),
        'Local': local_count,
        'International': intl_count
    })

df_bar = pd.DataFrame(bar_data)

if not df_bar.empty:
    try:
        fig, ax = plt.subplots(figsize=(12, 6))
        
        x = np.arange(len(df_bar))
        width = 0.35
        
        bars1 = ax.bar(x - width/2, df_bar['Local'], width, label='Local', color='#1f77b4')
        bars2 = ax.bar(x + width/2, df_bar['International'], width, label='International', color='#e377c2')
        
        ax.set_xlabel('Year', fontsize=12)
        ax.set_ylabel('Attendance Count', fontsize=12)
        ax.set_title('Overall: Local vs International Student Attendance by Year', fontsize=14, weight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(df_bar['Year'].astype(int))
        ax.legend()
        ax.grid(axis='y', alpha=0.3, linestyle='--')
        
        # Add value labels on bars
        ax.bar_label(bars1, padding=3, fontweight='bold')
        ax.bar_label(bars2, padding=3, fontweight='bold')
        
        plt.tight_layout()
        
        # Create overall summary table
        overall_type_counts = df_attended['Student_Type'].value_counts()
        df_overall_table = overall_type_counts.reset_index(name='Count').rename(columns={'index': 'Student Type'})
        df_overall_table['Percentage'] = (df_overall_table['Count'] / df_overall_table['Count'].sum() * 100).map('{:.1f}%'.format)
        
        figures_list.append({
            'fig': fig,
            'title': 'Overall Distribution (Yearly Comparison)',
            'table': df_overall_table
        })
    except Exception as e:
        pass

# 6. Summary Table (Year-by-Year Comparison)
type_counts = df_attended.groupby(['Workshop Timing_Year', 'Student_Type']).size().unstack(fill_value=0)

table_rows = []
for student_type in type_counts.columns:
    row_data = {
        "Student Type": student_type,
        "Total": type_counts[student_type].sum()
    }
    
    # Add yearly counts
    for year in years:
        row_data[str(year)] = type_counts.loc[year, student_type]
    
    # Add % Changes between consecutive years
    if len(years) >= 2:
        for i in range(len(years) - 1):
            y1, y2 = years[i], years[i+1]
            c1 = type_counts.loc[y1, student_type]
            c2 = type_counts.loc[y2, student_type]
            
            if c1 == 0:
                pct = "New" if c2 > 0 else "-"
            else:
                pct = f"{((c2 - c1) / c1) * 100:+.1f}%"
            
            row_data[f"% Change ({y1}->{y2})"] = pct
    
    table_rows.append(row_data)

df_table = pd.DataFrame(table_rows)
"""


# Q7 CODE
code_q7 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# 2. Logic: Sub-Category vs University
# Create/Clean University Column
if 'University Program' in df_attended.columns:
    df_attended['Uni_Clean'] = df_attended['University Program'].astype(str).apply(lambda x: x.split(',')[0]).str.title()
    uni_col = 'Uni_Clean'
    
    # Grouping Rule
    others_group = [
        'Grenoble Ecole De Management', 
        'La Trobe University', 
        'Monash College', 
        'The University Of Sydney'
    ]
    df_attended['Uni_Clean'] = df_attended['Uni_Clean'].apply(lambda x: 'Others' if x in others_group else x)
    
elif 'Uni_Clean' in df_attended.columns:
    uni_col = 'Uni_Clean'
else:
    uni_col = None

if uni_col:
    # Group by Sub-Category and University
    subcat_uni_counts = df_attended.groupby(['Sub-Category', uni_col]).size().unstack(fill_value=0)
    
    # 3. Stats
    kpi_result = {}
    figures_list = []
    
    # --- OVERALL ANALYSIS ---
    # 4a. Overall Graph - Top 10 Universities
    # Get top universities by total attendance
    uni_totals = subcat_uni_counts.sum(axis=0).sort_values(ascending=False).head(10)
    top_unis = uni_totals.index.tolist()
    
    # Filter to top universities
    subcat_uni_top = subcat_uni_counts[top_unis]
    
    fig, ax = plt.subplots(figsize=(14, 6))
    subcat_uni_top.plot(kind='bar', stacked=False, ax=ax, colormap='tab10')
    ax.set_title('Overall Workshop Attendance by Sub-Category & University (Top 10)')
    ax.set_xlabel('Sub-Category')
    ax.set_ylabel('Number of Attendees')
    ax.legend(title='University', bbox_to_anchor=(1.05, 1), loc='upper left')
    
    # Extend Y-axis by 15%
    if not subcat_uni_top.empty:
        ax.set_ylim(0, subcat_uni_top.max().max() * 1.15)
        
    ax.set_xticklabels([label.get_text().title() if isinstance(label.get_text(), str) else label.get_text() for label in ax.get_xticklabels()], rotation=45, ha='right')
    plt.tight_layout()
    
    # 4b. Overall Table
    df_overall = subcat_uni_top.copy()
    df_overall['Total'] = df_overall.sum(axis=1)
    df_overall = df_overall.sort_values('Total', ascending=False).reset_index()
    df_overall['Sub-Category'] = df_overall['Sub-Category'].apply(lambda x: x.title() if isinstance(x, str) else x)
    
    figures_list.append({
        'fig': fig,
        'title': 'Overall Analysis (Top 10 Universities)',
        'table': df_overall
    })
    
    # Use Overall table as default
    df_table = df_overall
    
    # --- YEARLY ANALYSIS ---
    years = sorted(df_attended['Workshop Timing_Year'].unique())
    for year in years:
        df_year = df_attended[df_attended['Workshop Timing_Year'] == year]
        if not df_year.empty:
            # Group by Sub-Category and University for this Year
            year_counts = df_year.groupby(['Sub-Category', uni_col]).size().unstack(fill_value=0)
            
            # Get top 10 universities for this year
            year_uni_totals = year_counts.sum(axis=0).sort_values(ascending=False).head(10)
            year_top_unis = year_uni_totals.index.tolist()
            year_counts_top = year_counts[year_top_unis]
            
            # Graph
            fig, ax = plt.subplots(figsize=(14, 6))
            year_counts_top.plot(kind='bar', stacked=False, ax=ax, colormap='tab10')
            ax.set_title(f'{year} Workshop Attendance by Sub-Category & University (Top 10)')
            ax.set_xlabel('Sub-Category')
            ax.set_ylabel('Number of Attendees')
            ax.legend(title='University', bbox_to_anchor=(1.05, 1), loc='upper left')
            
            # Extend Y-axis by 15%
            ax.set_ylim(0, year_counts_top.max().max() * 1.15)
            
            ax.set_xticklabels([label.get_text().title() if isinstance(label.get_text(), str) else label.get_text() for label in ax.get_xticklabels()], rotation=45, ha='right')
            plt.tight_layout()
            
            # Table
            df_year_table = year_counts_top.copy()
            df_year_table['Total'] = df_year_table.sum(axis=1)
            df_year_table = df_year_table.sort_values('Total', ascending=False).reset_index()
            df_year_table['Sub-Category'] = df_year_table['Sub-Category'].apply(lambda x: x.title() if isinstance(x, str) else x)
            
            figures_list.append({
                'fig': fig,
                'title': f'{year} Analysis (Top 10 Universities)',
                'table': df_year_table
            })
else:
    kpi_result = {'Status': f'University column not found. Available columns: {list(df_attended.columns)}'}
    figures_list = []
"""

# Q8 CODE
# Q8 CODE
# Q8 CODE
# Q8 CODE
code_q8 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# CLEAN UP Sub-Category (Title Case)
if 'Sub-Category' in df_attended.columns:
    df_attended['Sub-Category'] = df_attended['Sub-Category'].astype(str).str.title()
else:
    df_attended['Sub-Category'] = 'Uncategorized'

# 2. Identify Major Column
# Try 'Program Classification' -> 'Program' -> 'University Program'
if 'Program Classification' in df_attended.columns:
    major_col = 'Program Classification'
elif 'Program' in df_attended.columns:
    major_col = 'Program'
else:
    major_col = 'University Program'

# 3. Stats
kpi_result = {}
figures_list = []

# Validate Columns
if 'Sub-Category' not in df_attended.columns:
    kpi_result['Status'] = "Missing 'Sub-Category' column."
    df_table = pd.DataFrame([{"Error": kpi_result['Status']}])
elif major_col not in df_attended.columns:
    kpi_result['Status'] = f"Missing Major/Program column. (Checked for: {major_col})"
    df_table = pd.DataFrame([{"Error": kpi_result['Status']}])
else:
    # --- PER SUB-CATEGORY ANALYSIS ONLY ---
    # Get unique sub-categories
    unique_subcats = sorted(df_attended['Sub-Category'].dropna().unique())
    
    for subcat in unique_subcats:
        # Filter data for this sub-category (All Years)
        df_sub = df_attended[df_attended['Sub-Category'] == subcat]
        
        if not df_sub.empty:
            # Get Top 10 Majors for this Sub-Category
            major_counts_sub = df_sub[major_col].value_counts().head(10)
            
            if not major_counts_sub.empty:
                try:
                    fig, ax = plt.subplots(figsize=(10, 6))
                    
                    # Color palette
                    cmap = plt.get_cmap('tab10')
                    colors = [cmap(i) for i in range(len(major_counts_sub))]
                    
                    # Custom autopct function to hide % when < 1%
                    def autopct_format(pct):
                        return f'{pct:.1f}%' if pct >= 1 else ''
                    
                    wedges, texts, autotexts = ax.pie(
                        major_counts_sub.values,
                        labels=None,  # No labels on slices to avoid clutter
                        autopct=autopct_format,
                        startangle=90,
                        colors=colors,
                        pctdistance=0.85,
                        explode=[0.05 if i == 0 else 0 for i in range(len(major_counts_sub))]
                    )
                    
                    # Draw circle for donut style
                    centre_circle = plt.Circle((0,0),0.70,fc='white')
                    fig.gca().add_artist(centre_circle)
                    
                    ax.set_title(f'{subcat}: Major Classification Distribution', fontsize=14, weight='bold')
                    
                    # Legend
                    ax.legend(
                        major_counts_sub.index,
                        title=major_col,
                        loc="center left",
                        bbox_to_anchor=(1, 0, 0.5, 1)
                    )
                    
                    ax.axis('equal')
                    plt.tight_layout()
                    
                    # Table
                    df_sub_table = major_counts_sub.reset_index(name='Count').rename(columns={'index': major_col})
                    df_sub_table['Percentage'] = (df_sub_table['Count'] / df_sub_table['Count'].sum() * 100).map('{:.1f}%'.format)
                    
                    figures_list.append({
                        'fig': fig,
                        'title': f'{subcat} Distribution',
                        'table': df_sub_table
                    })
                except:
                    pass

    # Kpi
    kpi_result['Analysis Type'] = "Per Sub-Category Major Distribution (Pie Charts)"
    kpi_result['Sub-Categories Analyzed'] = len(unique_subcats)
    kpi_result['Data Scope'] = "All Data (No Yearly Split)"
"""

# Q9 CODE
# Q9 CODE
code_q9 = """
# 1. Filter
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean Year: Remove NaNs and convert to int
df_attended = df_attended.dropna(subset=['Workshop Timing_Year'])
df_attended['Workshop Timing_Year'] = df_attended['Workshop Timing_Year'].astype(int)

# 2. Logic: Calculate months until graduation using Grad attributes
def calculate_months_to_grad(row):
    try:
        # Get graduation year and month
        if pd.isna(row['Grad_Year']) or pd.isna(row['Grad_Month']):
            return pd.NA
        
        grad_year = int(row['Grad_Year'])
        
        # Convert month name to number
        month_map = {
            'january': 1, 'february': 2, 'march': 3, 'april': 4,
            'may': 5, 'june': 6, 'july': 7, 'august': 8,
            'september': 9, 'october': 10, 'november': 11, 'december': 12
        }
        grad_month_str = str(row['Grad_Month']).lower()
        grad_month = month_map.get(grad_month_str, pd.NA)
        
        if pd.isna(grad_month):
            return pd.NA
        
        # Get workshop date from Workshop Timing components
        if pd.isna(row['Workshop Timing_Year']) or pd.isna(row['Workshop Timing_Month']):
            return pd.NA
        
        workshop_year = int(row['Workshop Timing_Year'])
        workshop_month_str = str(row['Workshop Timing_Month']).lower()
        workshop_month = month_map.get(workshop_month_str, pd.NA)
        
        if pd.isna(workshop_month):
            return pd.NA
        
        # Calculate difference in months
        months_diff = (grad_year - workshop_year) * 12 + (grad_month - workshop_month)
        
        return months_diff
    except:
        return pd.NA

df_attended['Months_To_Grad'] = df_attended.apply(calculate_months_to_grad, axis=1)

# Bin into 4-month periods
def grad_quarter_bin(months):
    if pd.isna(months): 
        return 'Unknown'
    if months < 0:
        return 'Already Graduated'
    elif months <= 4:
        return '0-4 Months'
    elif months <= 8:
        return '5-8 Months'
    elif months <= 12:
        return '9-12 Months'
    elif months <= 16:
        return '13-16 Months'
    elif months <= 20:
        return '17-20 Months'
    elif months <= 24:
        return '21-24 Months'
    else:
        return '25+ Months'

df_attended['Grad_Quarter'] = df_attended['Months_To_Grad'].apply(grad_quarter_bin)

# 3. Stats
kpi_result = {}
figures_list = []

# Define order for quarters (4-month periods)
quarter_order = ['0-4 Months', '5-8 Months', '9-12 Months', '13-16 Months', '17-20 Months', '21-24 Months', '25+ Months', 'Already Graduated', 'Unknown']

# 4. Yearly Pie Charts
years = sorted(df_attended['Workshop Timing_Year'].unique())

for year in years:
    df_year = df_attended[df_attended['Workshop Timing_Year'] == year]
    
    if not df_year.empty:
        year_quarter_counts = df_year['Grad_Quarter'].value_counts()
        
        # Filter to existing quarters in this year
        year_quarter_counts = year_quarter_counts.reindex(
            [q for q in quarter_order if q in year_quarter_counts.index],
            fill_value=0
        )
        year_quarter_counts = year_quarter_counts[year_quarter_counts > 0]
        
        if not year_quarter_counts.empty:
            try:
                fig, ax = plt.subplots(figsize=(10, 6))
                
                # Color palette
                cmap = plt.get_cmap('tab10')
                colors = [cmap(i) for i in range(len(year_quarter_counts))]
                
                # Custom autopct function to hide % when < 1%
                def autopct_format(pct):
                    return f'{pct:.1f}%' if pct >= 1 else ''
                
                wedges, texts, autotexts = ax.pie(
                    year_quarter_counts.values,
                    labels=None,  # No labels on pie slices, use legend only
                    autopct=autopct_format,
                    startangle=90,
                    colors=colors,
                    pctdistance=0.85,
                    explode=[0.05 if i == 0 else 0 for i in range(len(year_quarter_counts))]
                )
                
                # Draw circle for donut style
                centre_circle = plt.Circle((0,0),0.70,fc='white')
                fig.gca().add_artist(centre_circle)
                
                ax.set_title(f'{int(year)}: Attendance by Time Until Graduation', fontsize=14, weight='bold')
                ax.legend(
                    year_quarter_counts.index,
                    title="Time to Graduation",
                    loc="center left",
                    bbox_to_anchor=(1, 0, 0.5, 1)
                )
                ax.axis('equal')
                plt.tight_layout()
                
                # Create table for this year
                df_year_table = year_quarter_counts.reset_index(name='Count').rename(columns={'index': 'Grad_Quarter'})
                df_year_table['Percentage'] = (df_year_table['Count'] / df_year_table['Count'].sum() * 100).map('{:.1f}%'.format)
                
                figures_list.append({
                    'fig': fig,
                    'title': f'{int(year)} Distribution',
                    'table': df_year_table
                })
            except Exception as e:
                pass

# 5. Overall Pie Chart
overall_quarter_counts = df_attended['Grad_Quarter'].value_counts()

# Filter to existing quarters overall
overall_quarter_counts = overall_quarter_counts.reindex(
    [q for q in quarter_order if q in overall_quarter_counts.index],
    fill_value=0
)
overall_quarter_counts = overall_quarter_counts[overall_quarter_counts > 0]

if not overall_quarter_counts.empty:
    try:
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Same color palette
        cmap = plt.get_cmap('tab10')
        colors = [cmap(i) for i in range(len(overall_quarter_counts))]
        
        # Custom autopct function to hide % when < 1%
        def autopct_format(pct):
            return f'{pct:.1f}%' if pct >= 1 else ''
        
        wedges, texts, autotexts = ax.pie(
            overall_quarter_counts.values,
            labels=None,  # No labels on pie slices, use legend only
            autopct=autopct_format,
            startangle=90,
            colors=colors,
            pctdistance=0.85,
            explode=[0.05 if i == 0 else 0 for i in range(len(overall_quarter_counts))]
        )
        
        # Draw circle for donut style
        centre_circle = plt.Circle((0,0),0.70,fc='white')
        fig.gca().add_artist(centre_circle)
        
        ax.set_title('Overall: Attendance by Time Until Graduation', fontsize=14, weight='bold')
        ax.legend(
            overall_quarter_counts.index,
            title="Time to Graduation",
            loc="center left",
            bbox_to_anchor=(1, 0, 0.5, 1)
        )
        ax.axis('equal')
        plt.tight_layout()
        
        # Create overall table
        df_overall_table = overall_quarter_counts.reset_index(name='Count').rename(columns={'index': 'Grad_Quarter'})
        df_overall_table['Percentage'] = (df_overall_table['Count'] / df_overall_table['Count'].sum() * 100).map('{:.1f}%'.format)
        
        figures_list.append({
            'fig': fig,
            'title': 'Overall Distribution',
            'table': df_overall_table
        })
    except Exception as e:
        pass

# 6. Summary Table - Cross-tabulation of quarters by year
df_table = df_attended.groupby(['Grad_Quarter', 'Workshop Timing_Year']).size().unstack(fill_value=0)

# Filter to existing quarters
existing_quarters = [q for q in quarter_order if q in df_table.index]
df_table = df_table.reindex(existing_quarters, fill_value=0)

# Get sorted years
years_in_table = sorted(df_table.columns)

# Add total column
df_table['Total'] = df_table.sum(axis=1)

# Add % Changes between consecutive years (before resetting index)
if len(years_in_table) >= 2:
    for i in range(len(years_in_table) - 1):
        y1, y2 = years_in_table[i], years_in_table[i+1]
        col_name = f'% Change ({y1}->{y2})'
        
        df_table[col_name] = df_table.apply(
            lambda row, start=y1, end=y2: (
                f"{((row[end] - row[start]) / row[start]) * 100:+.1f}%" 
                if row[start] != 0 
                else ("New" if row[end] > 0 else "-")
            ), 
            axis=1
        )

# Reset index for display
df_table = df_table.reset_index()
df_table.columns.name = None
"""

# Q10 CODE
# Q10 CODE
code_q10 = """
# 1. Setup
figures_list = []
kpi_result = {}

# We need ALL records to compare Attended vs Absent (Registered = All)
df_analysis = df.copy()

# 2. Logic: Reconstruct Date & Calculate Lead Days

# Clean Year
df_analysis = df_analysis.dropna(subset=['Workshop Timing_Year'])
df_analysis['Workshop Timing_Year'] = df_analysis['Workshop Timing_Year'].astype(int)

# Reconstruct Date
month_to_num = {
    'january': 1, 'february': 2, 'march': 3, 'april': 4,
    'may': 5, 'june': 6, 'july': 7, 'august': 8,
    'september': 9, 'october': 10, 'november': 11, 'december': 12
}
df_analysis['Month_Num'] = df_analysis['Workshop Timing_Month'].astype(str).str.lower().map(month_to_num)

# Filter valid dates for timing calculation
required_date_cols = ['Workshop Timing_Year', 'Month_Num', 'Workshop Timing_DayNumber', 'Registered Date']
df_analysis = df_analysis.dropna(subset=required_date_cols)

if df_analysis.empty:
    kpi_result = {'Status': 'No available data with complete timing info'}
    df_table = pd.DataFrame([{"Error": kpi_result['Status']}])
else:
    # Build Date
    df_analysis['Attended_Date_Reconstructed'] = pd.to_datetime(
        df_analysis[['Workshop Timing_Year', 'Month_Num', 'Workshop Timing_DayNumber']].rename(
            columns={'Workshop Timing_Year': 'year', 'Month_Num': 'month', 'Workshop Timing_DayNumber': 'day'}
        ),
        errors='coerce'
    )
    df_analysis = df_analysis.dropna(subset=['Attended_Date_Reconstructed'])
    
    if df_analysis.empty:
        kpi_result = {'Status': 'Date reconstruction failed'}
        df_table = pd.DataFrame([{"Error": kpi_result['Status']}])
    else:
        # Calculate Lead Days
        df_analysis['Lead_Days'] = (df_analysis['Attended_Date_Reconstructed'] - df_analysis['Registered Date']).dt.days
        
        # Binning
        def lead_bin(days):
            if pd.isna(days): return 'Unknown'
            if days <= 0: return 'Same Day'
            elif days <= 3: return '1-3 Days'
            elif days <= 7: return '1 Week'
            else: return '>1 Week'
        
        df_analysis['Reg_Timing'] = df_analysis['Lead_Days'].apply(lead_bin)
        
        # Define Status: Attended vs Absent
        # If 'Attended', status is 'Attended'. All else is 'Absent'.
        df_analysis['Status_Simple'] = df_analysis['Attendance Status'].astype(str).str.lower().apply(
            lambda x: 'Attended' if x == 'attended' else 'Absent'
        )
        
        # Timing Order
        timing_order = ['Same Day', '1-3 Days', '1 Week', '>1 Week']
        
        # --- Helper for plotting ---
        def plot_grouped_bar(data_df, title_prefix, timing_order, color_registered='#00d2d3', color_attended='#2e86de'):
            # Group by Timing and Status
            grouped = data_df.groupby(['Reg_Timing', 'Status_Simple']).size().unstack(fill_value=0)
            
            # Ensure both columns exist (Attended, Absent)
            for col in ['Attended', 'Absent']:
                if col not in grouped.columns:
                    grouped[col] = 0
            
            # Create 'Registered' (Total) column
            grouped['Registered'] = grouped['Attended'] + grouped['Absent']
            
            # Reindex rows to timing order
            valid_timings = [t for t in timing_order if t in grouped.index]
            grouped = grouped.reindex(valid_timings)
            
            if grouped.empty:
                return None, None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Plot Registered vs Attended
            x = np.arange(len(grouped))
            width = 0.35
            
            rects1 = ax.bar(x - width/2, grouped['Registered'], width, label='Registered', color=color_registered)
            rects2 = ax.bar(x + width/2, grouped['Attended'], width, label='Attended', color=color_attended)
            
            ax.set_ylabel('Count')
            ax.set_title(f'{title_prefix}: Registered vs Attended by Registration Timing', fontsize=14, weight='bold')
            ax.set_xticks(x)
            ax.set_xticklabels(grouped.index, rotation=0)
            
            # Extend Y-axis by 15%
            max_val = grouped[['Registered', 'Attended']].max().max()
            ax.set_ylim(0, max_val * 1.15)
            
            ax.legend()
            
            # Add labels
            def autolabel(rects):
                for rect in rects:
                    height = rect.get_height()
                    if height > 0:
                        ax.annotate(f'{height}',
                                    xy=(rect.get_x() + rect.get_width() / 2, height),
                                    xytext=(0, 3),
                                    textcoords="offset points",
                                    ha='center', va='bottom')
            
            autolabel(rects1)
            autolabel(rects2)
            
            plt.tight_layout()
            
            # Create Table
            df_t = grouped[['Registered', 'Attended']].copy()
            df_t['Attendance Rate'] = (df_t['Attended'] / df_t['Registered'] * 100).map('{:.1f}%'.format)
            df_t = df_t.reset_index().rename(columns={'Reg_Timing': 'Registration Timing'})
            
            return fig, df_t
            
        
        # --- 3. Overall Analysis ---
        try:
            fig_all, table_all = plot_grouped_bar(df_analysis, 'Overall', timing_order)
            if fig_all:
                figures_list.append({
                    'fig': fig_all,
                    'title': 'Overall Analysis',
                    'table': table_all
                })
                # Set table_all as main df_table for display
                df_table = table_all
        except Exception as e:
            df_table = pd.DataFrame([{"Error": f"Overall plot error: {str(e)}"}])

        # --- 4. Yearly Analysis ---
        years = sorted(df_analysis['Workshop Timing_Year'].unique())
        for year in years:
            df_year = df_analysis[df_analysis['Workshop Timing_Year'] == year]
            if not df_year.empty:
                try:
                    fig_y, table_y = plot_grouped_bar(df_year, f'{year}', timing_order)
                    if fig_y:
                        figures_list.append({
                            'fig': fig_y,
                            'title': f'{year} Analysis',
                            'table': table_y
                        })
                except Exception as e:
                    pass

        # KPI Stats
        total_recs = len(df_analysis)
        total_attended = len(df_analysis[df_analysis['Status_Simple'] == 'Attended'])
        kpi_result = {
            'Total Registered': total_recs,
            'Total Attended': total_attended,
            'Overall Attendance Rate': f"{(total_attended/total_recs*100):.1f}%"
        }
"""



# Q11 CODE (Formerly Q12)
code_q11 = """

# 1. Setup
figures_list = []
kpi_result = {}

# 2. Logic: Unique Students & Top List per University
# Filter Attended
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean University
df_attended['Uni_Clean'] = df_attended['University Program'].astype(str).apply(lambda x: x.split(',')[0]).str.title()

# Grouping Rule
others_group = [
    'Grenoble Ecole De Management', 
    'La Trobe University', 
    'Monash College', 
    'The University Of Sydney'
]
df_attended['Uni_Clean'] = df_attended['Uni_Clean'].apply(lambda x: 'Others' if x in others_group else x)

# Clean Name
if 'Student Name' not in df.columns:
    df_attended['Display_Name'] = df_attended['SIMID']
else:
    df_attended['Display_Name'] = df_attended['Student Name'].fillna(df_attended['SIMID'])

# Group by University - Unique Count of SIMID
uni_unique_counts = df_attended.groupby('Uni_Clean')['SIMID'].nunique().sort_values(ascending=False)

# Top 10 Students per University
top_students_data = []

# Iterate over each university present
for uni in uni_unique_counts.index:
    # Filter for this uni
    df_uni = df_attended[df_attended['Uni_Clean'] == uni]
    
    # Count attendance per student
    student_counts = df_uni['Display_Name'].value_counts()
    
    # Get Top 10
    top_10 = student_counts.head(10)
    
    # Format list: "Name (Count)"
    top_list_str = ", ".join([f"{name} ({count})" for name, count in top_10.items()])
    
    top_students_data.append({
        'University': uni,
        'Unique Students': uni_unique_counts[uni],
        'Top 10 High Attendance Students': top_list_str
    })

df_table = pd.DataFrame(top_students_data)

# 3. Visualization: Bar Count of Unique Students
if not uni_unique_counts.empty:
    fig, ax = plt.subplots(figsize=(12, 8))
    
    # Plot
    top_n_plot = uni_unique_counts.head(15) # Show top 15 for plotting clarity if many
    
    sns.barplot(x=top_n_plot.values, y=top_n_plot.index, palette='viridis', ax=ax)
    
    ax.set_title('Top Universities by Unique Student Participation', fontsize=14, weight='bold')
    ax.set_xlabel('Unique Student Count')
    ax.set_ylabel('University')
    
    # Extend X-axis by 15% (for horizontal bar plot, X is count)
    if not top_n_plot.empty:
        ax.set_xlim(0, top_n_plot.max() * 1.15)
    
    # Add numbers
    for i, v in enumerate(top_n_plot.values):
        ax.text(v + 0.1, i, str(v), color='black', va='center', fontweight='bold')
        
    plt.tight_layout()
    
    figures_list.append({
        'fig': fig,
        'title': 'Unique Participants by University',
        'table': df_table
    })
else:
    kpi_result['Status'] = 'No attendance data found.'
    fig = None

# 4. KPI Stats
kpi_result['Total Unique Students'] = df_attended['SIMID'].nunique()
if not uni_unique_counts.empty:
    kpi_result['Top University'] = f"{uni_unique_counts.index[0]} ({uni_unique_counts.iloc[0]})"
"""


# Q12 CODE (Formerly Q13)
code_q12 = """
# 1. Setup
figures_list = []
kpi_result = {}

# 2. Logic: Monthly Attendance by University
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Clean University
df_attended['Uni_Clean'] = df_attended['University Program'].astype(str).apply(lambda x: x.split(',')[0]).str.title()

# Apply Grouping Rule
# Group specific universities into "Others"
others_group = [
    'Grenoble Ecole De Management', 
    'La Trobe University', 
    'Monash College', 
    'The University Of Sydney'
]

df_attended['Uni_Grouped'] = df_attended['Uni_Clean'].apply(
    lambda x, g=others_group: 'Others' if x in g else x
)

# Date Construction (YYYY-MM)
month_to_num = {
    'january': 1, 'february': 2, 'march': 3, 'april': 4,
    'may': 5, 'june': 6, 'july': 7, 'august': 8,
    'september': 9, 'october': 10, 'november': 11, 'december': 12
}

# Ensure columns exist
if 'Workshop Timing_Year' in df_attended.columns and 'Workshop Timing_Month' in df_attended.columns:
    df_attended = df_attended.dropna(subset=['Workshop Timing_Year', 'Workshop Timing_Month'])
    df_attended['Year'] = df_attended['Workshop Timing_Year'].astype(int)
    df_attended['Month_Num'] = df_attended['Workshop Timing_Month'].astype(str).str.lower().map(month_to_num)
    
    # Filter valid months
    df_attended = df_attended.dropna(subset=['Month_Num'])
    df_attended['Month_Num'] = df_attended['Month_Num'].astype(int)
    
    # Create Sortable Key and Label
    df_attended['Time_Key'] = df_attended['Year'] * 100 + df_attended['Month_Num']
    df_attended['Month_Label'] = df_attended['Year'].astype(str) + '-' + df_attended['Month_Num'].astype(str).str.zfill(2)
    
    # Aggregate
    # Group by Time Key (for sort), Month Label (for display), and University
    monthly_counts = df_attended.groupby(['Time_Key', 'Month_Label', 'Uni_Grouped']).size().reset_index(name='Attendance Count')
    
    # Sort chronologically
    monthly_counts = monthly_counts.sort_values('Time_Key')
    
    # Pivot for Table
    df_pivot = monthly_counts.pivot_table(
        index=['Time_Key', 'Month_Label'], 
        columns='Uni_Grouped', 
        values='Attendance Count', 
        fill_value=0
    ).reset_index().sort_values('Time_Key')
    
    # Calculate Percentages and Format
    uni_cols = [c for c in df_pivot.columns if c not in ['Time_Key', 'Month_Label']]
    
    # Calculate Row Totals
    df_pivot['Row_Total'] = df_pivot[uni_cols].sum(axis=1)
    
    df_table = df_pivot.copy()
    
    for col in uni_cols:
        # Format: Count (Percent%)
        # Capture col as c to avoid scope issues in exec loop
        df_table[col] = df_pivot.apply(
            lambda row, c=col: f"{int(row[c])} ({row[c]/row['Row_Total']*100:.1f}%)" if row['Row_Total'] > 0 else f"{int(row[c])}", 
            axis=1
        )
            
    # Final cleanup
    df_table = df_table.drop(columns=['Time_Key', 'Row_Total'])
    
    # 3. Visualization: Line Graph
    if not monthly_counts.empty:
        fig, ax = plt.subplots(figsize=(14, 8))
        
        # We need to identify Top N universities to visualize cleanly, plus Others
        # or just visualize all from the Grouped set (which already consolidated some)
        
        # Get list of unis to plot (Sort by total attendance to put biggest first in legend)
        uni_totals = monthly_counts.groupby('Uni_Grouped')['Attendance Count'].sum().sort_values(ascending=False)
        ranked_unis = uni_totals.index.tolist()
        
        # Plot
        sns.lineplot(
            data=monthly_counts, 
            x='Month_Label', 
            y='Attendance Count', 
            hue='Uni_Grouped', 
            hue_order=ranked_unis,
            marker='o',
            linewidth=2,
            ax=ax
        )
        
        ax.set_title('Monthly Attendance Trends by University', fontsize=14, weight='bold')
        ax.set_xlabel('Month')
        ax.set_ylabel('Attendance Count')
        plt.xticks(rotation=45)
        ax.grid(True, alpha=0.3, linestyle='--')
        
        # Extend Y-axis by 15%
        if not monthly_counts.empty:
            ax.set_ylim(0, monthly_counts['Attendance Count'].max() * 1.15)
            
        ax.legend(title='University', bbox_to_anchor=(1.02, 1), loc='upper left')
        
        plt.tight_layout()
        
        figures_list.append({
            'fig': fig,
            'title': 'Monthly Trends by University',
            'table': df_table
        })
        
        # 4. KPI Stats
        kpi_result['Total Attended'] = len(df_attended)
        if not uni_totals.empty:
            kpi_result['Top University'] = uni_totals.index[0]
            
    else:
        kpi_result['Status'] = 'No valid monthly data found.'
        fig = None

else:
    kpi_result['Status'] = 'Missing Year/Month columns.'
    fig = None
"""


# Q13 CODE
code_q13 = """
# 1. Setup
figures_list = []
kpi_result = {}

# 2. Logic: Trainer by Attendance
# Filter Attended
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Ensure Trainer column exists
if 'Trainer' in df_attended.columns:
    # Clean Trainer Name (Title Case, Strip)
    df_attended['Trainer'] = df_attended['Trainer'].fillna('Unknown').astype(str).str.strip().str.title()
    
    # Filter out empty or 'Nan'
    df_attended = df_attended[~df_attended['Trainer'].isin(['', 'Nan', 'None', 'Unknown'])]
    
    # Count by Trainer
    trainer_counts = df_attended['Trainer'].value_counts().reset_index()
    trainer_counts.columns = ['Trainer', 'Attendance Count']
    
    # Sort
    trainer_counts = trainer_counts.sort_values('Attendance Count', ascending=True) # Ascending for barh
    
    # Calculate Percentage
    total_attended = trainer_counts['Attendance Count'].sum()
    trainer_counts['Percentage'] = (trainer_counts['Attendance Count'] / total_attended * 100).map('{:.1f}%'.format)
    
    # 3. Table
    # Sort descending for table
    df_table = trainer_counts.sort_values('Attendance Count', ascending=False).reset_index(drop=True)
    
    # 4. Visualization: Horizontal Bar Chart
    if not trainer_counts.empty:
        fig, ax = plt.subplots(figsize=(10, max(6, len(trainer_counts) * 0.4)))
        
        # Plot
        bars = ax.barh(trainer_counts['Trainer'], trainer_counts['Attendance Count'], color='#1f77b4')
        
        ax.set_title('Attendance by Trainer', fontsize=14, weight='bold')
        ax.set_xlabel('Attendance Count', fontsize=12)
        ax.set_ylabel('Trainer', fontsize=12)
        ax.grid(axis='x', alpha=0.3, linestyle='--')
        
        # Extend X-axis by 15% (Horizontal Bar)
        if not trainer_counts.empty:
            ax.set_xlim(0, trainer_counts['Attendance Count'].max() * 1.15)
        
        # Add labels
        ax.bar_label(bars, padding=3, fontweight='bold')
        
        plt.tight_layout()
        
        figures_list.append({
            'fig': fig,
            'title': 'Attendance by Trainer',
            'table': df_table
        })
        
        # 5. KPI Stats
        kpi_result['Total Trainers'] = len(trainer_counts)
        if not trainer_counts.empty:
            kpi_result['Top Trainer'] = f"{trainer_counts.iloc[-1]['Trainer']} ({trainer_counts.iloc[-1]['Attendance Count']})"
            
    else:
        kpi_result['Status'] = 'No valid trainer data found.'
        fig = None

else:
    kpi_result['Status'] = "Column 'Trainer' not found in dataset."
    fig = None
"""

# Q14 CODE
code_q14 = """
# 1. Setup
figures_list = []
kpi_result = {}

# 2. Logic: Workshop Title by Attendance (Table Only)
# Filter Attended
df_attended = df[df['Attendance Status'].astype(str).str.lower() == 'attended'].copy()

# Ensure Columns exist
if 'Event Name' in df.columns:
    # --- A. Overall Statistics ---
    # Group by Event Name
    # Calculate Attendance Count and Unique Runs (based on Attended Date)
    # If Attended Date is missing, Runs defaults to 0 or 1? 
    # Let's count uniques.
    
    aggs = {'Attendance Status': 'count'}
    if 'Attended Date' in df_attended.columns:
        # Count unique dates as Runs
        # Note: If multiple sessions on same day, this counts as 1 run. 
        # Ideally we use Session ID but that's not guaranteed. Date is best proxy.
        overall_stats = df_attended.groupby('Event Name').agg(
            Attendance_Count=('Attendance Status', 'count'),
            Runs=('Attended Date', 'nunique')
        ).reset_index()
    else:
        overall_stats = df_attended.groupby('Event Name').agg(
            Attendance_Count=('Attendance Status', 'count')
        ).reset_index()
        overall_stats['Runs'] = "N/A"

    overall_stats.columns = ['Workshop Title', 'Attendance Count', 'No. of Runs']
    
    # Sort Descending by Attendance
    overall_stats = overall_stats.sort_values('Attendance Count', ascending=False).reset_index(drop=True)
    
    # Add Overall to figures
    figures_list.append({
        'title': 'Overall Workshop Statistics',
        'table': overall_stats,
        'fig': None
    })
    
    # --- B. Yearly Breakdown ---
    if 'Workshop Timing_Year' in df_attended.columns:
        # Group by Year and Event Name
        if 'Attended Date' in df_attended.columns:
            yearly_stats = df_attended.groupby(['Workshop Timing_Year', 'Event Name']).agg(
                Attendance_Count=('Attendance Status', 'count'),
                Runs=('Attended Date', 'nunique')
            ).reset_index()
        else:
            yearly_stats = df_attended.groupby(['Workshop Timing_Year', 'Event Name']).agg(
                Attendance_Count=('Attendance Status', 'count')
            ).reset_index()
            yearly_stats['Runs'] = "N/A"
            
        yearly_stats.columns = ['Year', 'Workshop Title', 'Attendance Count', 'No. of Runs']
        
        # Sort by Year (desc) then Attendance (desc)
        yearly_stats = yearly_stats.sort_values(['Year', 'Attendance Count'], ascending=[False, False]).reset_index(drop=True)
        
        # Add Yearly to figures
        figures_list.append({
            'title': 'Yearly Workshop Breakdown',
            'table': yearly_stats,
            'fig': None
        })
    
    # 5. KPI Stats
    kpi_result['Total Workshops'] = len(overall_stats)
    if not overall_stats.empty:
        kpi_result['Top Workshop'] = f"{overall_stats.iloc[0]['Workshop Title']} ({overall_stats.iloc[0]['Attendance Count']})"
    else:
        kpi_result['Status'] = 'No attendance data found.'

else:
    kpi_result['Status'] = "Column 'Event Name' not found."
"""


# ==============================================================================
# RENDER ALL BLOCKS
# ==============================================================================

# ==============================================================================
# 3-PAGE WORKFLOW & PPT GENERATION
# ==============================================================================

# --- HELPER: EXECUTE ANALYSIS CODE (Decoupled from Display) ---
def execute_analysis_code(q_id, df_to_use, code_to_run, exclude_uni=False):
    """
    Executes the analysis code and returns the results dictionary.
    Does NOT allow st. calls (sandboxed).
    
    Args:
        q_id: Question ID
        df_to_use: DataFrame to analyze
        code_to_run: Python code string to execute
        exclude_uni: If True, exclude university-specific events from analysis
    """
    # Apply exclusion filter if requested
    if exclude_uni:
        # Filter out events containing any exclusion term (case-insensitive)
        pattern = '|'.join([re.escape(term) for term in EXCLUSION_TERMS])
        # Ensure 'Event Name' exists
        if 'Event Name' in df_to_use.columns:
            mask = df_to_use['Event Name'].astype(str).str.contains(pattern, case=False, na=False)
            df_to_use = df_to_use[~mask].copy()
    
    local_vars = {
        'df': df_to_use, 
        'pd': pd, 
        'plt': plt, 
        'sns': sns,
        'np': __import__('numpy'),
        'kpi_result': {},
        'fig': None,
        'df_table': None,
        'figures_list': []
    }
    
    try:
        exec(code_to_run, globals(), local_vars)
        return {
            'success': True,
            'kpi_result': local_vars.get('kpi_result', {}),
            'figures_list': local_vars.get('figures_list', []),
            'fig': local_vars.get('fig'),
            'df_table': local_vars.get('df_table'),
            'error': None
        }
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }

# --- HELPER: GENERATE POWERPOINT ---
def generate_ppt(df_global, exclude_uni=False):
    """
    Generates a PowerPoint presentation with all 10 analysis questions.
    
    Args:
        df_global: The global dataframe with all data
        exclude_uni: If True, exclude university-specific events from analysis
    
    Returns a BytesIO object containing the PPTX file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        st.error("python-pptx library not found. Please install it.")
        return None

    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Workshop Attendance Analysis"
    subtitle.text = f"Generated on {pd.Timestamp.now().strftime('%Y-%m-%d')}"

    # Codes Map
    codes_map = {
        1: code_q1, 2: code_q2, 3: code_q3, 4: code_q4, 5: code_q5,
        6: code_q6, 7: code_q7, 8: code_q8, 9: code_q9, 10: code_q10,
        11: code_q11, 12: code_q12, 13: code_q13, 14: code_q14
    }
    
    titles = [
        "Attendance and Attrition",
        "Unique Participant Count",
        "Most Popular Days & Time Slots",
        "Attendance by University",
        "Workshop Attendance by Sub-Category",
        "Attendance by Student Type (Local vs International)",
        "Workshop Attendance by Sub-Category & Student Type",
        "Workshop Attendance by Sub-Category & Academic Major",
        "Attendance by Expected Graduation Period",
        "Registered vs Attended by Registration Timing",
        "Unique Counts & Top Students by University",
        "Monthly Attendance by University",
        "Trainer by Attendance",
        "Workshop Titles by Attendance"
    ]

    for q_id in range(1, 15):
        q_title = titles[q_id-1]
        code = st.session_state.get(f"edited_code_{q_id}", codes_map[q_id])
        
        # Execute Analysis with exclude_uni setting
        result = execute_analysis_code(q_id, df_global, code, exclude_uni=exclude_uni)
        
        # Add Question Title Slide (Section Divider)
        title_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Add title text box
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)
        
        txBox = title_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"Q{q_id}: {q_title}"
        
        # Format the title
        p = tf.paragraphs[0]
        p.alignment = 1  # Center alignment
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        
        if not result['success']:
            # Error Slide
            slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
            slide.shapes.title.text = f"Q{q_id}: {q_title}"
            slide.placeholders[1].text = f"Error generating analysis: {result['error']}"
            continue

        items_to_plot = []
        
        # Collect all items (Single or Multi)
        if result['figures_list']:
            for item in result['figures_list']:
                items_to_plot.append({
                    'title': item.get('title', f"Q{q_id} Analysis"),
                    'fig': item.get('fig'),
                    'table': item.get('table')
                })
        elif result['fig'] or (result['df_table'] is not None and not result['df_table'].empty):
            items_to_plot.append({
                'title': f"Q{q_id}: {q_title}",
                'fig': result['fig'],
                'table': result['df_table']
            })

        if not items_to_plot:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Q{q_id}: {q_title}"
            slide.placeholders[1].text = "No data available for this analysis."
            continue

        # Create Slides for each item
        for item in items_to_plot:
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            
            # Add Title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item['title']
            p.font.size = Pt(24)
            p.font.bold = True

            # Add Image (Left or Center)
            if item['fig']:
                image_stream = io.BytesIO()
                try:
                    item['fig'].savefig(image_stream, format='png', dpi=150, bbox_inches='tight')
                    image_stream.seek(0)
                    slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1.5), height=Inches(4.5))
                except Exception as e:
                    pass

            # Add Table (Right) - Simplified text version for now as pptx tables are complex
            # Ideally we would iterate and build a real pptx table
            if item['table'] is not None and not item['table'].empty:
                df = item['table']
                # Limit rows to prevent overflow
                df_show = df.head(15) 
                
                rows, cols = df_show.shape
                # Position table to the right of the image
                left = Inches(5.5) if item['fig'] else Inches(1)
                top = Inches(1.5)
                width = Inches(4.0) if item['fig'] else Inches(8.0)
                height = Inches(0.8) # minimal height
                
                shape = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
                
                # Function to set cell borders
                from pptx.oxml.xmlchemy import OxmlElement
                
                def SubElement(parent, tagname, **kwargs):
                    element = OxmlElement(tagname)
                    element.attrib.update(kwargs)
                    parent.append(element)
                    return element

                def _set_cell_border(cell, border_color="000000", border_width='12700'): # 12700 = 1pt
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
                        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
                        solidFill = SubElement(ln, 'a:solidFill')
                        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
                        prstDash = SubElement(ln, 'a:prstDash', val='solid')
                        round_ = SubElement(ln, 'a:round')
                        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
                        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

                # Column headers
                for i, col_name in enumerate(df_show.columns):
                    cell = shape.cell(0, i)
                    cell.text = str(col_name)
                    cell.fill.background() # No fill
                    _set_cell_border(cell, border_color="000000", border_width='12700')  # Black borders
                    # Bold headers with black text
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
                
                # Rows
                for r in range(rows):
                    for c in range(cols):
                        cell = shape.cell(r+1, c)
                        cell.text = str(df_show.iloc[r, c])
                        cell.fill.background() # No fill
                        _set_cell_border(cell, border_color="000000", border_width='12700')  # Black borders
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
                
                # Add note if truncated
                if len(df) > 15:
                    txBox = slide.shapes.add_textbox(left, top + Inches(4.5), width, Inches(0.5))
                    txBox.text_frame.text = f"(Showing top 15 of {len(df)} rows)"

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ==============================================================================
# MAIN APP LOGIC
# ==============================================================================

def main():
    # 1. Page Config
    # st.set_page_config ... (already set globally)
    
    # --- GLOBAL HEADER ---
    # 2. Session State for Navigation
    if 'current_page' not in st.session_state:
        st.session_state['current_page'] = 1
    
    # Helper to switch pages
    def go_to_page(page_num):
        st.session_state['current_page'] = page_num
        st.rerun()

    # --- HEADER & NAVIGATION ---
    # Layout: Left Button | Title | Right Button
    col_nav_1, col_nav_2, col_nav_3 = st.columns([1, 6, 1])

    with col_nav_2:
        st.markdown("<h1 style='text-align: center; margin-bottom: 0px;'>Career Development Analysis Tool</h1>", unsafe_allow_html=True)
        if st.session_state.get('current_page', 1) == 1:
             st.markdown("<p style='text-align: center; margin-top: 0px;'>Upload attendance logs and taxonomy to generate the dashboard</p>", unsafe_allow_html=True)

    with col_nav_1:
        if st.session_state['current_page'] == 2:
            if st.button("⬅ Data Prep", key="nav_back_2", use_container_width=True):
                go_to_page(1)
        elif st.session_state['current_page'] == 3:
             if st.button("⬅ Analysis", key="nav_back_3", use_container_width=True):
                go_to_page(2)

    with col_nav_3:
        if st.session_state['current_page'] == 2:
             if st.button("Report ➡", key="nav_next_2", use_container_width=True):
                go_to_page(3)

    # ==============================================================================
    # PAGE 1: DATA PREPARATION
    # ==============================================================================
    if st.session_state['current_page'] == 1:
        # Title removed as requested
        pass
        
        run_data_ingestion()
        
        # --- QUESTIONS OVERVIEW SECTION ---
        st.markdown("---")
        st.markdown("### 📊 Available Analysis Questions")
        st.markdown("Once your data is processed, you'll have access to the following analysis questions:")
        
        # Create a nice grid layout for questions
        questions_data = [
            ("Q1: Attendance and Attrition", QUESTION_PURPOSES[1]),
            ("Q2: Unique Participant Count", QUESTION_PURPOSES[2]),
            ("Q3: Most Popular Days & Time Slots", QUESTION_PURPOSES[3]),
            ("Q4: Attendance by University", QUESTION_PURPOSES[4]),
            ("Q5: Workshop Attendance by Sub-Category", QUESTION_PURPOSES[5]),
            ("Q6: Attendance by Student Type", QUESTION_PURPOSES[6]),
            ("Q7: Workshop Attendance by Sub-Category & University", QUESTION_PURPOSES[7]),
            ("Q8: Workshop Attendance by Sub-Category & Academic Major", QUESTION_PURPOSES[8]),
            ("Q9: Attendance by Expected Graduation Period", QUESTION_PURPOSES[9]),
            ("Q10: Registered vs Attended by Registration Timing", QUESTION_PURPOSES[10]),
            ("Q11: Unique Counts & Top Students by University", QUESTION_PURPOSES[11]),
            ("Q12: Monthly Attendance by University", QUESTION_PURPOSES[12]),
            ("Q13: Trainer by Attendance", QUESTION_PURPOSES[13]),
            ("Q14: Workshop Titles by Attendance", QUESTION_PURPOSES[14])
        ]
        
        # Display in 2 columns
        col1, col2 = st.columns(2)
        
        # Calculate split point (round up for left column)
        mid_point = (len(questions_data) + 1) // 2
        
        for i, (title, description) in enumerate(questions_data):
            # i < mid_point goes to col1 (Left), rest to col2 (Right)
            with col1 if i < mid_point else col2:
                st.markdown(f"**{title}**")
                st.markdown(f"{description}")
                st.markdown("")  # Add spacing



    # ==============================================================================
    # PAGE 2: ANALYSIS (Q1-Q10)
    # ==============================================================================
    elif st.session_state['current_page'] == 2:
        
        titles = [
            "Attendance and Attrition",
            "Unique Participant Count",
            "Most Popular Days & Time Slots",
            "Attendance by University",
            "Workshop Attendance by Sub-Category",
            "Attendance by Student Type (Local vs International)",
            "Workshop Attendance by Sub-Category & University",
            "Workshop Attendance by Sub-Category & Academic Major",
            "Attendance by Expected Graduation Period",
            "Registered vs Attended by Registration Timing",
            "Unique Counts & Top Students by University",
            "Monthly Attendance by University",
            "Trainer by Attendance",
            "Workshop Titles by Attendance"
        ]
        
        # Code Map
        default_codes = [code_q1, code_q2, code_q3, code_q4, code_q5, code_q6, code_q7, code_q8, code_q9, code_q10, code_q11, code_q12, code_q13, code_q14]
        
        # Render all questions
        for i in range(14):
            q_id = i + 1
            render_sandbox(q_id, titles[i], default_codes[i])
            

    # ==============================================================================
    # PAGE 3: REPORT GENERATION
    # ==============================================================================
    elif st.session_state['current_page'] == 3:
        # Title removed as requested
        st.markdown("<h5 style='text-align: center;'>Generate a PowerPoint presentation summarizing the analysis from all 14 questions.</h5>", unsafe_allow_html=True)
        
        # Exclude Uni Events Option
        st.markdown("")
        col_opt1, col_opt2, col_opt3 = st.columns([1, 2, 1])
        with col_opt2:
            exclude_uni_ppt = st.checkbox(
                "🎓 Exclude University-Specific Events",
                value=False,
                key="exclude_uni_ppt",
                help="Exclude workshops designed specifically for particular universities from the PowerPoint report"
            )
        
        st.markdown("")
        _, col_center, _ = st.columns([1, 1, 1])
        
        with col_center:
            if st.button("Generate PowerPoint Presentation", type="primary", use_container_width=True):
                with st.spinner("Generating slides..."):
                    # Generate PPT with exclude uni events setting
                    ppt_io = generate_ppt(st.session_state['data'], exclude_uni=exclude_uni_ppt)
                    
                    if ppt_io:
                        st.success("✅ Presentation Generated!")
                        
                        # Preview (using a placeholder concept or just download)
                        st.info("Click below to download your report.")
                        
                        st.download_button(
                            label="📥 Download PowerPoint (.pptx)",
                            data=ppt_io,
                            file_name=f"Workshop_Analysis_Report_{pd.Timestamp.now().strftime('%Y%m%d_%H%M')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )




    # ==============================================================================
    # FOOTER
    # ==============================================================================
    st.markdown("---")
    st.markdown(
        """
        <div style='text-align: center; color: #666; font-size: 12px;'>
            @Copyright 2025 Learner's Advisory and Career Connect, SIM
        </div>
        """,
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()
